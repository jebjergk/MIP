"""
Build MIP Behind the Scenes Presentation
How MIP actually works — the mechanics, the math, the gates — explained simply.

Companion to the Platform Introduction deck. Focused on *what goes on under the hood*
rather than what you see on screen. Uses analogies and visual illustrations throughout
to make complex mechanics accessible to a normal person.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── MIP palette ──────────────────────────────────────────────────────────────
MIP_BLUE      = RGBColor(48, 105, 240)
MIP_DARK      = RGBColor(15, 23, 42)
MIP_SLATE     = RGBColor(51, 65, 85)
MIP_BG        = RGBColor(245, 248, 252)
MIP_TEXT      = RGBColor(31, 41, 55)
MIP_MUTED     = RGBColor(95, 107, 122)
MIP_GREEN     = RGBColor(46, 125, 50)
MIP_AMBER     = RGBColor(146, 95, 0)
MIP_RED       = RGBColor(198, 40, 40)
MIP_PURPLE    = RGBColor(106, 27, 154)
MIP_TEAL      = RGBColor(0, 121, 107)
MIP_CYAN      = RGBColor(0, 151, 167)
MIP_INDIGO    = RGBColor(57, 73, 171)
WHITE         = RGBColor(255, 255, 255)
LIGHT_GRAY    = RGBColor(215, 220, 228)
CARD_BG       = RGBColor(255, 255, 255)
DARK_CARD_BG  = RGBColor(15, 23, 42)
DARK_CARD_BDR = RGBColor(36, 48, 65)
HERO_GREEN_BG = RGBColor(230, 245, 234)
HERO_BLUE_BG  = RGBColor(227, 236, 253)
HERO_PURP_BG  = RGBColor(243, 229, 245)
HERO_AMB_BG   = RGBColor(255, 248, 225)
HERO_RED_BG   = RGBColor(254, 235, 235)
FOOTER_BG     = RGBColor(238, 242, 248)
FOOTER_TEXT    = RGBColor(110, 123, 141)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
#  SHARED HELPERS (same design system as Platform Intro deck)
# ══════════════════════════════════════════════════════════════════════════════

def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _draw_header(slide, title, subtitle, section="BEHIND THE SCENES"):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.78))
    top.fill.solid()
    top.fill.fore_color.rgb = MIP_DARK
    top.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.75), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()
    sec = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.14), Inches(2.3), Inches(0.42))
    sec.fill.solid()
    sec.fill.fore_color.rgb = RGBColor(36, 48, 71)
    sec.line.color.rgb = MIP_BLUE
    tf = sec.text_frame
    tf.text = section
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(225, 233, 246)
    t = slide.shapes.add_textbox(Inches(0.7), Inches(1.02), Inches(9.5), Inches(0.7))
    tf = t.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(34)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MIP_DARK
    s = slide.shapes.add_textbox(Inches(0.7), Inches(1.72), Inches(11.8), Inches(0.55))
    tf = s.text_frame
    tf.word_wrap = True
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = MIP_MUTED


def _add_card(slide, x, y, w, h, title, items, edge=MIP_BLUE, title_size=15, item_size=12):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = LIGHT_GRAY
    card.shadow.inherit = False
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = edge
    bar.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(title_size)
    p0.font.color.rgb = MIP_SLATE
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(item_size)
        p.font.color.rgb = MIP_TEXT
        p.space_before = Pt(4)


def _add_chip(slide, x, y, w, text, bg, fg=WHITE, h=0.34):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    chip.fill.solid()
    chip.fill.fore_color.rgb = bg
    chip.line.fill.background()
    tf = chip.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = fg


def _add_process(slide, labels, y, start_x=0.75, box_w=1.88, gap=0.22):
    for idx, label in enumerate(labels):
        x = start_x + idx * (box_w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = RGBColor(207, 214, 225)
        tf = box.text_frame
        tf.text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = MIP_SLATE
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + box_w + 0.02), Inches(y + 0.2), Inches(0.18), Inches(0.45))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MIP_BLUE
            arrow.line.fill.background()


def _add_analogy_bar(slide, icon, text, y=6.6):
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(240, 243, 250)
    bar.line.color.rgb = LIGHT_GRAY
    tf = bar.text_frame
    tf.word_wrap = True
    tf.text = f"{icon}  {text}"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = MIP_SLATE


def _add_footer(slide, text="MIP Behind the Scenes  |  How the engine actually works"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), SLIDE_W, Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOOTER_BG
    bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.6), Inches(7.22), Inches(10.0), Inches(0.2))
    tf = t.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = FOOTER_TEXT


def _add_hero_box(slide, x, y, w, h, bg, title, body_lines, title_color=MIP_DARK):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = LIGHT_GRAY
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = title_color
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = MIP_TEXT
        p.space_before = Pt(4)


def _add_stat_box(slide, x, y, w, h, value, label, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = LIGHT_GRAY
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = value
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(26)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    p = tf.add_paragraph()
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.color.rgb = MIP_MUTED
    p.space_before = Pt(2)


def _add_formula_box(slide, x, y, w, h, formula, explanation, color=MIP_BLUE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(240, 243, 255)
    box.line.color.rgb = color
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = formula
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    p = tf.add_paragraph()
    p.text = explanation
    p.font.size = Pt(11)
    p.font.color.rgb = MIP_TEXT
    p.space_before = Pt(4)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def _slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_DARK)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = MIP_DARK
    slide.shapes[-1].line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.0), Inches(0.15), Inches(5.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(1.3), Inches(1.2), Inches(8.0), Inches(1.0))
    tf = t.text_frame
    tf.text = "MIP \u2014 Behind the Scenes"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    s = slide.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(8.0), Inches(0.6))
    tf = s.text_frame
    tf.text = "How the engine actually works \u2014 from raw market data to live trades"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)

    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(3.1), Inches(3.0), Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = MIP_BLUE
    div.line.fill.background()

    desc = slide.shapes.add_textbox(Inches(1.3), Inches(3.4), Inches(6.5), Inches(2.0))
    tf = desc.text_frame
    tf.word_wrap = True
    lines = [
        "This deck explains what happens under the hood:",
        "the calculations, the gates, the decision logic.",
        "",
        "No jargon without explanation.",
        "No formula without an analogy.",
    ]
    tf.text = lines[0]
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = RGBColor(203, 213, 225)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_before = Pt(4)

    pillars = [
        ("SIGNALS", MIP_BLUE), ("TRAINING", MIP_TEAL),
        ("GATES", MIP_PURPLE), ("EXECUTION", MIP_GREEN),
    ]
    for i, (label, color) in enumerate(pillars):
        _add_chip(slide, 1.3 + i * 2.2, 6.1, 1.9, label, color)


def _slide_big_picture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Big Picture", "MIP runs a continuous loop: observe the market, learn from it, decide, act, and learn from the outcome", "OVERVIEW")

    _add_process(slide, ["Ingest", "Detect", "Evaluate", "Train", "Propose", "Gate", "Execute", "Learn"], y=2.1, start_x=0.4, box_w=1.35, gap=0.18)

    _add_hero_box(slide, 0.7, 3.3, 5.9, 3.2, HERO_BLUE_BG, "\U0001F52C  Research Loop (runs daily)", [
        "1. Ingest OHLCV bars from AlphaVantage / IBKR",
        "2. Detect patterns (momentum, ORB, mean reversion, pullbacks)",
        "3. Evaluate: did past signals hit or miss after 1, 3, 5, 10, 20 days?",
        "4. Train: update maturity scores and trust labels per symbol",
        "5. The system gets smarter every day it runs",
    ], title_color=MIP_BLUE)

    _add_hero_box(slide, 6.75, 3.3, 5.9, 3.2, HERO_GREEN_BG, "\U0001F3AF  Decision Loop (runs daily)", [
        "6. Propose: pick the best signals for each portfolio",
        "7. Gate: trust gate, risk gate, capacity gate, freshness gate",
        "8. Committee: AI reviews each proposal (can block)",
        "9. Execute: place in simulation or live via IBKR",
        "10. Learn: record outcome, feed it back into training",
    ], title_color=MIP_GREEN)

    _add_analogy_bar(slide, "\U0001F504", "Like a chef who tastes every dish (evaluate), adjusts the recipe (train), and only serves when quality passes inspection (gate).")
    _add_footer(slide)


def _slide_how_signals_are_born(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "How Signals Are Born", "MIP watches price movements and fires a signal when a pattern's conditions are met", "SIGNALS")

    _add_process(slide, ["OHLCV Bars", "Returns", "Pattern Check", "Signal Fired"], y=2.1, start_x=1.0, box_w=2.5, gap=0.5)

    _add_card(slide, 0.7, 3.3, 4.0, 3.2, "Momentum Pattern (Daily)", [
        "Calculates simple return: (Close - PrevClose) / PrevClose",
        "Counts recent positive days (e.g. 3 of last 20)",
        "Checks if price is at a new high within window",
        "Calculates z-score: how unusual is today's move?",
        "Signal fires when ALL conditions pass simultaneously",
    ], edge=MIP_BLUE)

    _add_card(slide, 4.85, 3.3, 4.0, 3.2, "Intraday Patterns (15-min bars)", [
        "ORB: First bar defines the range; breakout above or below",
        "Mean Reversion: Price drifts far from average, then snaps back",
        "Pullback: Strong move, brief pause, then continuation",
        "Each pattern has its own parameter set (configurable)",
        "One signal per symbol per session (no duplicates)",
    ], edge=MIP_PURPLE)

    _add_card(slide, 9.0, 3.3, 3.65, 3.2, "What Gets Recorded", [
        "Symbol + pattern ID",
        "Timestamp of the signal",
        "Score (strength measure)",
        "Market type (FX, stock, ETF)",
        "Direction (bullish / bearish)",
    ], edge=MIP_TEAL)

    _add_analogy_bar(slide, "\U0001F6A8", "Like a smoke detector: it watches continuously, has specific trigger conditions, and fires only when ALL thresholds are crossed \u2014 not just one.")
    _add_footer(slide)


def _slide_horizons(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Horizons: Looking Into the Future", "After a signal fires, MIP checks what happened next \u2014 at 5 different time distances", "EVALUATION")

    _add_stat_box(slide, 0.7, 2.15, 2.2, 1.2, "H1", "1 Day Ahead", MIP_BLUE)
    _add_stat_box(slide, 3.05, 2.15, 2.2, 1.2, "H3", "3 Days Ahead", MIP_TEAL)
    _add_stat_box(slide, 5.4, 2.15, 2.2, 1.2, "H5", "5 Days Ahead", MIP_GREEN)
    _add_stat_box(slide, 7.75, 2.15, 2.2, 1.2, "H10", "10 Days Ahead", MIP_PURPLE)
    _add_stat_box(slide, 10.1, 2.15, 2.5, 1.2, "H20", "20 Days Ahead", MIP_INDIGO)

    _add_formula_box(slide, 0.7, 3.6, 5.9, 0.8,
        "Realized Return = (Future Price \u2212 Signal Price) / Signal Price",
        "If price was $100 at signal and $103 after 5 days: H5 return = +3.0%")

    _add_card(slide, 0.7, 4.65, 5.9, 1.9, "HIT, MISS, or NEUTRAL?", [
        "HIT:  return \u2265 +0.2%  (the signal was right \u2014 price went up)",
        "MISS:  return \u2264 -0.2%  (the signal was wrong \u2014 price dropped)",
        "NEUTRAL:  in between  (inconclusive \u2014 nothing happened)",
        "These labels are counted to build the HIT RATE per pattern",
    ], edge=MIP_GREEN, item_size=11)

    _add_card(slide, 6.75, 3.6, 5.9, 2.95, "Why 5 Horizons?", [
        "A pattern might work over 1 day but fail over 20",
        "Or it might be meaningless short-term but powerful long-term",
        "By measuring at 5 distances, MIP finds the sweet spot",
        "Training builds a profile: \"Pattern X works best at H5 for EUR/USD\"",
        "The system doesn't guess \u2014 it measures and remembers",
        "",
        "Intraday uses different horizons:",
        "  H1 = 1 bar (15 min), H4 = 1 hour, H8 = 2 hours, EOD = end of day",
    ], edge=MIP_PURPLE, item_size=11)

    _add_analogy_bar(slide, "\U0001F321\uFE0F", "Like a weather forecast: predicting tomorrow is different from predicting next week. MIP measures accuracy at each distance to know which forecasts to trust.")
    _add_footer(slide)


def _slide_training_lab(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Training Lab", "How MIP measures whether a pattern has earned enough evidence to be trusted", "TRAINING")

    _add_formula_box(slide, 0.7, 2.1, 11.9, 0.85,
        "Maturity Score = Sample Score (max 30) + Coverage Score (max 40) + Horizon Score (max 30)  =  0 to 100",
        "Each component is measured independently and capped. A pattern must excel in all three to reach CONFIDENT status.")

    _add_stat_box(slide, 0.7, 3.2, 3.8, 1.5, "30 pts", "Sample Size\nDo we have enough signals?", MIP_BLUE)
    _add_stat_box(slide, 4.65, 3.2, 3.8, 1.5, "40 pts", "Coverage\nWhat % of outcomes evaluated?", MIP_GREEN)
    _add_stat_box(slide, 8.5, 3.2, 4.1, 1.5, "30 pts", "Horizon Depth\nHow many of 5 horizons tested?", MIP_PURPLE)

    _add_chip(slide, 0.7, 5.0, 2.5, "< 25 INSUFFICIENT", MIP_RED)
    _add_chip(slide, 3.3, 5.0, 2.5, "25-49 WARMING UP", MIP_AMBER)
    _add_chip(slide, 5.9, 5.0, 2.5, "50-74 LEARNING", MIP_BLUE)
    _add_chip(slide, 8.5, 5.0, 2.5, "\u2265 75 CONFIDENT", MIP_GREEN)

    _add_card(slide, 0.7, 5.6, 11.9, 0.9, "What This Means in Practice", [
        "A pattern with 10 signals, 50% coverage, and 2 horizons tested would score about 35 \u2014 WARMING UP. Not allowed to trade yet.",
        "A pattern with 40+ signals, 90% coverage, and all 5 horizons would score 90+ \u2014 CONFIDENT. Ready for proposals.",
    ], edge=MIP_TEAL, item_size=11)

    _add_analogy_bar(slide, "\U0001F393", "Like a student's transcript: sample size is how many exams taken, coverage is how many subjects, horizons is depth of knowledge. All three matter for graduation.")
    _add_footer(slide)


def _slide_trust_gate(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Trust Gate", "Maturity alone isn't enough \u2014 the pattern must also prove it actually makes money recently", "TRUST")

    _add_card(slide, 0.7, 2.1, 6.0, 2.0, "Pattern-Level Trust (Global)", [
        "\u2265 30 evaluated outcomes AND \u2265 80% coverage rate",
        "Average return > 0 (historically profitable)",
        "Recent 90 days: \u2265 10 signals, \u2265 50% hit rate, positive avg return",
        "Only patterns meeting ALL criteria earn TRUSTED status",
    ], edge=MIP_GREEN)

    _add_card(slide, 6.85, 2.1, 5.8, 2.0, "Symbol-Level Trust (Local)", [
        "\u2265 8 recent signals for this specific symbol",
        "Recent hit rate \u2265 50% for this symbol",
        "Recent average return \u2265 0% for this symbol",
        "Even a globally trusted pattern fails if it doesn't work for THIS symbol",
    ], edge=MIP_PURPLE)

    _add_chip(slide, 0.7, 4.35, 2.0, "TRUSTED", MIP_GREEN)
    _add_chip(slide, 2.8, 4.35, 1.7, "WATCH", MIP_AMBER)
    _add_chip(slide, 4.6, 4.35, 2.2, "UNTRUSTED", MIP_RED)
    _add_chip(slide, 6.9, 4.35, 2.5, "LOW EVIDENCE", RGBColor(113, 128, 150))

    _add_card(slide, 0.7, 4.95, 5.9, 1.55, "What Happens If Trust Fails?", [
        "If a pattern is not TRUSTED: no proposals are generated for it",
        "If a symbol's local trust is not TRUSTED: BUY proposals are blocked",
        "The system keeps learning \u2014 trust can be earned over time",
        "Trust can also be LOST if recent performance deteriorates",
    ], edge=MIP_RED)

    _add_card(slide, 6.75, 4.95, 5.9, 1.55, "Why Two Layers of Trust?", [
        "Global trust: \"This pattern works in general across many symbols\"",
        "Local trust: \"This pattern works specifically for EUR/USD\"",
        "A pattern could work for stocks but not for FX pairs",
        "Double-checking prevents false confidence from aggregation",
    ], edge=MIP_TEAL)

    _add_analogy_bar(slide, "\U0001F3E6", "Like a bank loan: your overall credit score matters (global), but so does your history with THIS specific bank (local). Both must pass.")
    _add_footer(slide)


def _slide_committee_room(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Committee Room", "Every trade proposal must survive a gauntlet of 5 independent checks before it can execute", "GATES")

    gates = [
        ("1. Trust Gate", "Is the pattern trusted for this symbol?\nBUY blocked if symbol trust \u2260 TRUSTED", MIP_GREEN),
        ("2. Risk Gate", "Is the portfolio in drawdown-stop mode?\nBlocked if drawdown \u2265 stop threshold", MIP_RED),
        ("3. Capacity Gate", "Are there slots available?\nBlocked if positions \u2265 max allowed", MIP_AMBER),
        ("4. Freshness Gate", "Is market data fresh?\nBlocked if daily bars are stale (>30 hrs old)", MIP_BLUE),
        ("5. AI Committee", "Does the AI agree?\nClaude reviews and can block the trade", MIP_PURPLE),
    ]

    y = 2.1
    for title, desc, color in gates:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(0.82))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = LIGHT_GRAY
        edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(0.08), Inches(0.82))
        edge.fill.solid()
        edge.fill.fore_color.rgb = color
        edge.line.fill.background()

        gate_chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.4), Inches(y + 0.2), Inches(1.0), Inches(0.38))
        gate_chip.fill.solid()
        gate_chip.fill.fore_color.rgb = color
        gate_chip.line.fill.background()
        tf = gate_chip.text_frame
        tf.text = "GATE"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        tf = box.text_frame
        tf.word_wrap = True
        parts = desc.split("\n")
        tf.text = f"  {title}"
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        for part in parts:
            p = tf.add_paragraph()
            p.text = f"      {part}"
            p.font.size = Pt(11)
            p.font.color.rgb = MIP_TEXT
            p.space_before = Pt(1)

        y += 0.92

    _add_analogy_bar(slide, "\U0001F6A7", "Like airport security: ID check (trust), bag scan (risk), seat availability (capacity), boarding pass time (freshness), and a final officer review (AI).")
    _add_footer(slide)


def _slide_position_sizing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Position Sizing & Costs", "How MIP decides HOW MUCH to buy, and what the real cost of a trade looks like", "MECHANICS")

    _add_formula_box(slide, 0.7, 2.1, 5.9, 0.85,
        "Target Value = Cash \u00D7 Max Position %",
        "Example: $100,000 cash \u00D7 5% = $5,000 per position. Never bet more than 5% on one trade.")

    _add_formula_box(slide, 6.75, 2.1, 5.9, 0.85,
        "Quantity = Target Value / Entry Price",
        "Example: $5,000 / $150 per share = 33 shares. Rounded to whole shares.")

    _add_card(slide, 0.7, 3.2, 4.0, 2.0, "Entry Price Adjustment", [
        "You don't get the exact close price",
        "Slippage: market moves while order fills",
        "Spread: difference between bid and ask",
        "Buy Price = Close \u00D7 (1 + slippage + spread/2)",
        "Sell Price = Close \u00D7 (1 \u2212 slippage \u2212 spread/2)",
    ], edge=MIP_BLUE, item_size=11)

    _add_card(slide, 4.85, 3.2, 4.0, 2.0, "Trading Fees", [
        "Fee = max(Minimum Fee, Notional \u00D7 Fee Rate)",
        "Example: max($1, $5,000 \u00D7 0.001) = $5",
        "Fees are subtracted from P&L on both entry and exit",
        "Total cost: buy fee + sell fee + slippage + spread",
    ], edge=MIP_AMBER, item_size=11)

    _add_card(slide, 9.0, 3.2, 3.65, 2.0, "Committee Can Scale", [
        "AI committee can reduce size via SIZE_FACTOR (0\u20131)",
        "If SIZE_FACTOR = 0.5: position halved",
        "Lower conviction = smaller bet",
        "Committee can also block entirely",
    ], edge=MIP_PURPLE, item_size=11)

    _add_card(slide, 0.7, 5.45, 11.9, 1.05, "Market Type Quota", [
        "MIP splits available slots: 60% for stocks, 40% for FX. This prevents over-concentration in one asset class.",
        "If there are 10 open slots: up to 6 go to stocks, up to 4 go to FX. Unused slots backfill from the other pool.",
    ], edge=MIP_TEAL, item_size=11)

    _add_analogy_bar(slide, "\U0001F3B0", "Like a poker player who never puts more than 5% of their chips on one hand \u2014 and adjusts the bet based on how confident they are in their cards.")
    _add_footer(slide)


def _slide_simulation_engine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Simulation Engine", "How MIP replays the market day-by-day, opening and closing positions to test strategy quality", "SIMULATION")

    _add_process(slide, ["Start Cash", "Daily Bars", "Check Exits", "Check Entries", "Update Equity"], y=2.1, start_x=0.5, box_w=2.2, gap=0.3)

    _add_card(slide, 0.7, 3.3, 4.0, 3.2, "How Positions Are Managed", [
        "Entry: buy at adjusted close price (with slippage)",
        "Hold period: determined by the signal's horizon",
        "Exit: auto-sell after N bars (e.g. H5 = 5 days)",
        "No manual stop-loss/take-profit in simulation",
        "FIFO ordering: first position in, first out",
        "Holidays carry forward \u2014 uses last known price",
    ], edge=MIP_BLUE)

    _add_card(slide, 4.85, 3.3, 4.0, 3.2, "Daily Equity Tracking", [
        "Each day: recalculate all position values",
        "Total Equity = Cash + Open Position Values",
        "Daily P&L = Today's Equity \u2212 Yesterday's Equity",
        "Daily Return = Daily P&L / Yesterday's Equity",
        "Peak Equity = highest equity ever reached",
        "Drawdown = (Peak \u2212 Current) / Peak",
    ], edge=MIP_GREEN)

    _add_card(slide, 9.0, 3.3, 3.65, 3.2, "Risk Brakes", [
        "Drawdown Stop: if drawdown hits threshold",
        "(e.g. 10%), new entries are BLOCKED",
        "Bust Protection: if equity falls below",
        "bust level, all entries stop",
        "Cooldown: forced waiting period",
        "after risk event before re-entering",
    ], edge=MIP_RED)

    _add_analogy_bar(slide, "\u2708\uFE0F", "Like a flight simulator: pilots train by \"flying\" through recorded weather patterns. MIP trades through recorded market days to learn what works before risking real money.")
    _add_footer(slide)


def _slide_news_scoring(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "News Intelligence: The Math", "How MIP converts raw headlines into a quantified, bounded influence score", "NEWS")

    _add_formula_box(slide, 0.7, 2.1, 5.9, 0.85,
        "Z-Score = (Today's News Count \u2212 7-Day Average) / 7-Day StdDev",
        "Measures how unusual today's news volume is compared to the recent norm. Z \u2265 1.5 = HOT.")

    _add_formula_box(slide, 6.75, 2.1, 5.9, 0.85,
        "Novelty Score = 1 \u2212 (Duplicate Count / Total Count)",
        "How much of today's news is genuinely new vs. recycled. High novelty = more attention-worthy.")

    _add_card(slide, 0.7, 3.2, 4.0, 2.1, "Burst Detection", [
        "Track 7-day rolling average and std deviation",
        "Today's count significantly above average = burst",
        "Z \u2265 1.5 \u2192 HOT badge (unusual activity)",
        "Z < 1.5 \u2192 NORMAL badge",
        "No news at all \u2192 NONE badge",
    ], edge=MIP_RED)

    _add_card(slide, 4.85, 3.2, 4.0, 2.1, "Uncertainty Detection", [
        "Scans headlines for both bullish AND bearish keywords",
        "If BOTH directions appear: UNCERTAINTY = true",
        "Mixed signals mean the market is confused",
        "Uncertainty flag adds caution to proposals",
    ], edge=MIP_AMBER)

    _add_card(slide, 9.0, 3.2, 3.65, 2.1, "Bounded Influence", [
        "News score adjustment capped at \u00B120%",
        "Currently diagnostic only \u2014 committee",
        "sees the data but ranking is unchanged",
        "No trade is ever placed purely on news",
    ], edge=MIP_TEAL)

    _add_card(slide, 0.7, 5.55, 11.9, 0.95, "The Key Principle: News Informs, Never Overrides", [
        "News can make the committee more cautious or more confident, but it cannot bypass trust gates, risk gates, or capacity limits.",
        "The influence is mathematically bounded (\u00B120% cap) to prevent a single headline from distorting the entire system.",
    ], edge=MIP_PURPLE, item_size=12)

    _add_analogy_bar(slide, "\U0001F4F0", "Like a judge reading the morning paper: it may inform their thinking, but the verdict is based on evidence presented in court \u2014 not on headlines.")
    _add_footer(slide)


def _slide_parallel_worlds(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Parallel Worlds: The What-If Machine", "MIP asks: 'What would have happened if we'd used a different strategy?' \u2014 and calculates the answer", "COUNTERFACTUAL")

    scenarios = [
        ("BASELINE", "What if we did nothing and stayed in cash?", "\U0001F4B5"),
        ("SIZING", "What if each position was 2\u00D7 or 0.5\u00D7 the size?", "\U0001F4CF"),
        ("TIMING", "What if we waited 1\u20135 bars before entering?", "\u23F1\uFE0F"),
        ("THRESHOLD", "What if we used stricter/looser signal filters?", "\U0001F50D"),
        ("HORIZON", "What if we held positions longer or shorter?", "\U0001F4C5"),
        ("EARLY EXIT", "What if we took profits earlier at a multiplier?", "\U0001F3C3"),
    ]

    y = 2.1
    for i, (name, desc, icon) in enumerate(scenarios):
        col = i % 3
        row = i // 3
        sx = 0.7 + col * 4.2
        sy = y + row * 1.55
        _add_card(slide, sx, sy, 4.0, 1.35, f"{icon}  {name}", [desc], edge=MIP_BLUE if row == 0 else MIP_PURPLE, item_size=11)

    _add_formula_box(slide, 0.7, 5.4, 5.9, 0.7,
        "Regret = Best Alternative Return \u2212 Actual Return",
        "If cash returned +2% and you returned +5%, regret = \u22123% (no regret!)")

    _add_card(slide, 6.75, 5.4, 5.9, 0.7, "Policy Health", [
        "If actual outperforms all alternatives \u2192 HEALTHY.  If one alternative consistently wins \u2192 investigate.",
    ], edge=MIP_GREEN, item_size=12)

    _add_analogy_bar(slide, "\U0001F30D", "Like testing alternate routes on Google Maps AFTER your trip: 'Would the highway have been faster?' MIP answers this for every trading day.")
    _add_footer(slide)


def _slide_live_bridge(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Live Bridge to IBKR", "How an approved trade travels from MIP's brain to Interactive Brokers' execution engine", "EXECUTION")

    _add_process(slide, ["Validated", "Committee", "PM Accept", "Compliance", "Revalidate", "Execute"], y=2.1, start_x=0.35, box_w=1.85, gap=0.25)

    _add_card(slide, 0.7, 3.3, 4.0, 1.7, "Pre-Flight Checks", [
        "Price guard: current price vs. signal price",
        "Exposure check: would this exceed limits?",
        "Market open: is the exchange actually trading?",
        "Symbol halt: is this stock suspended?",
    ], edge=MIP_BLUE)

    _add_card(slide, 4.85, 3.3, 4.0, 1.7, "Revalidation (just before execution)", [
        "Fresh 1-minute bar required (not stale)",
        "Price deviation check: has price moved too far?",
        "Exposure recalculated with current NAV",
        "Must pass ALL checks again \u2014 market moves fast",
    ], edge=MIP_AMBER)

    _add_card(slide, 9.0, 3.3, 3.65, 1.7, "Order Structure", [
        "Parent: Market or Limit order",
        "Child 1: Take-profit (limit order)",
        "Child 2: Stop-loss (stop order)",
        "Bracket order: all three linked",
    ], edge=MIP_GREEN)

    _add_card(slide, 0.7, 5.25, 5.9, 1.25, "Broker Snapshot Sync", [
        "MIP pulls from IBKR: NAV, cash, positions, open orders, executions",
        "Written to BROKER_SNAPSHOTS table every sync cycle",
        "Drift detection: does MIP's view match the broker's reality?",
    ], edge=MIP_TEAL)

    _add_card(slide, 6.75, 5.25, 5.9, 1.25, "Risk Brakes (Live)", [
        "Drawdown stop: if portfolio drops \u2265 threshold from peak \u2192 block entries",
        "Remains blocked until all positions are flat (closed)",
        "Only then can the system resume trading \u2014 automatic cooldown",
    ], edge=MIP_RED)

    _add_analogy_bar(slide, "\U0001F680", "Like a rocket launch: pre-flight checks, countdown hold points, a final 'go' from mission control, and real-time telemetry monitoring once airborne.")
    _add_footer(slide)


def _slide_daily_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Daily Pipeline", "Every weekday at 5 PM ET, MIP runs a 10-step automated pipeline \u2014 here's what happens", "PIPELINE")

    steps = [
        ("1", "Ingest Market Data", "Pull latest OHLCV bars from AlphaVantage / IBKR", MIP_BLUE),
        ("2", "Generate Signals", "Run pattern detection on fresh bars \u2192 fire new signals", MIP_BLUE),
        ("3", "Evaluate Outcomes", "Check if yesterday's signals hit, missed, or were neutral", MIP_TEAL),
        ("4", "Run Backtests", "Update hit rates, avg returns, and pattern scores", MIP_TEAL),
        ("5", "Train Patterns", "Recalculate maturity, trust labels, activate/deactivate", MIP_TEAL),
        ("6", "Run Simulation", "Simulate portfolio: open/close positions, track equity", MIP_GREEN),
        ("7", "Propose Trades", "Select best trusted signals \u2192 create proposals", MIP_PURPLE),
        ("8", "Validate & Execute", "Run through all gates + committee \u2192 execute approved", MIP_PURPLE),
        ("9", "Generate Briefs", "AI creates daily digest + training narratives", MIP_INDIGO),
        ("10", "Parallel Worlds", "Run counterfactual scenarios + regret analysis", MIP_INDIGO),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = 0.7 + col * 6.35
        y = 2.1 + row * 0.88

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.15), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = LIGHT_GRAY
        edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(0.75))
        edge.fill.solid()
        edge.fill.fore_color.rgb = color
        edge.line.fill.background()

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.2), Inches(y + 0.15), Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        tf = circle.text_frame
        tf.text = num
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

        tf = box.text_frame
        tf.word_wrap = True
        tf.text = f"       {title}"
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = MIP_SLATE
        p = tf.add_paragraph()
        p.text = f"          {desc}"
        p.font.size = Pt(10)
        p.font.color.rgb = MIP_MUTED
        p.space_before = Pt(1)

    _add_analogy_bar(slide, "\U0001F3ED", "Like a factory assembly line that runs every evening: raw materials in, quality checks at every station, finished product out. Every step is logged.")
    _add_footer(slide)


def _slide_learning_loop(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "The Learning Loop", "How every outcome feeds back into the system to make tomorrow's decisions better", "LEARNING")

    _add_process(slide, ["Signal", "Proposal", "Trade", "Outcome", "Update Training"], y=2.1, start_x=0.5, box_w=2.2, gap=0.3)

    arrow_back = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.7), Inches(3.15), Inches(11.9), Inches(0.25))
    arrow_back.fill.solid()
    arrow_back.fill.fore_color.rgb = RGBColor(220, 228, 240)
    arrow_back.line.fill.background()
    arrow_back.rotation = 180.0
    lbl = slide.shapes.add_textbox(Inches(4.5), Inches(3.15), Inches(4.5), Inches(0.25))
    tf = lbl.text_frame
    tf.text = "Feedback loop \u2014 outcomes update training data"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = MIP_MUTED

    _add_card(slide, 0.7, 3.65, 4.0, 2.8, "What Gets Updated?", [
        "Hit rate: % of signals that were right",
        "Average return: mean profit/loss per signal",
        "Cumulative return: total historical edge",
        "Pattern score: hit_rate \u00D7 cumulative_return",
        "Maturity stage: may advance or regress",
        "Trust label: may upgrade or downgrade",
        "IS_ACTIVE flag: can switch on or off",
    ], edge=MIP_BLUE)

    _add_card(slide, 4.85, 3.65, 4.0, 2.8, "What Gets Recorded?", [
        "Learning Decision Ledger captures everything:",
        "  BEFORE state vs AFTER state",
        "  Influence delta (what caused the change)",
        "  Causality links to upstream events",
        "  Outcome state when the position closes",
        "",
        "3 event types: TRAINING, DECISION, LIVE",
    ], edge=MIP_PURPLE)

    _add_card(slide, 9.0, 3.65, 3.65, 2.8, "Virtuous Cycle", [
        "Good patterns: earn higher scores, get",
        "  more proposals, build more evidence",
        "Bad patterns: lose trust, get fewer",
        "  proposals, eventually deactivated",
        "The system self-corrects over time",
        "",
        "No manual tuning required",
    ], edge=MIP_GREEN)

    _add_analogy_bar(slide, "\U0001F3C8", "Like a sports team reviewing game film: every play is analyzed, every mistake becomes a lesson, and the playbook is updated before the next game.")
    _add_footer(slide)


def _slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_DARK)

    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = MIP_DARK
    slide.shapes[-1].line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.0), Inches(0.15), Inches(5.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(1.3), Inches(1.2), Inches(10.5), Inches(1.2))
    tf = t.text_frame
    tf.text = "Under the Hood, It's All Evidence"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(241, 245, 249)

    s = slide.shapes.add_textbox(Inches(1.3), Inches(2.5), Inches(10.0), Inches(0.6))
    tf = s.text_frame
    tf.word_wrap = True
    tf.text = "MIP doesn't predict the market. It measures what works, earns trust through evidence, and only acts when the math supports it."
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)

    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(3.3), Inches(4.0), Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = MIP_BLUE
    div.line.fill.background()

    points = [
        ("\U0001F6A8", "Signals", "Patterns fire based on math, not intuition"),
        ("\U0001F393", "Training", "Maturity is earned through evidence, not declared"),
        ("\U0001F6A7", "Gates", "5 independent checks must ALL pass before execution"),
        ("\U0001F4CA", "Sizing", "Position size is calculated, never guessed"),
        ("\U0001F4F0", "News", "Influence is bounded \u2014 headlines can't override gates"),
        ("\U0001F30D", "Parallel Worlds", "Every strategy is stress-tested against alternatives"),
        ("\U0001F504", "Learning", "Every outcome updates the system's knowledge"),
    ]

    y = 3.55
    for icon, label, desc in points:
        row = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(10.0), Inches(0.42))
        tf = row.text_frame
        tf.word_wrap = True
        tf.text = f"{icon}  {label}  \u2014  {desc}"
        tf.paragraphs[0].font.size = Pt(15)
        tf.paragraphs[0].font.color.rgb = RGBColor(203, 213, 225)
        y += 0.47


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD DECK
# ══════════════════════════════════════════════════════════════════════════════

def build_deck(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs)               # 1  - Title
    _slide_big_picture(prs)         # 2  - The Big Picture
    _slide_how_signals_are_born(prs)  # 3  - How Signals Are Born
    _slide_horizons(prs)            # 4  - Horizons
    _slide_training_lab(prs)        # 5  - The Training Lab
    _slide_trust_gate(prs)          # 6  - The Trust Gate
    _slide_committee_room(prs)      # 7  - The Committee Room
    _slide_position_sizing(prs)     # 8  - Position Sizing & Costs
    _slide_simulation_engine(prs)   # 9  - The Simulation Engine
    _slide_news_scoring(prs)        # 10 - News Intelligence: The Math
    _slide_parallel_worlds(prs)     # 11 - Parallel Worlds
    _slide_live_bridge(prs)         # 12 - The Live Bridge to IBKR
    _slide_daily_pipeline(prs)      # 13 - The Daily Pipeline
    _slide_learning_loop(prs)       # 14 - The Learning Loop
    _slide_closing(prs)             # 15 - Closing

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    deck_path = Path(__file__).with_name("MIP_Behind_The_Scenes.pptx")
    build_deck(deck_path)
    print(f"Created: {deck_path}")
