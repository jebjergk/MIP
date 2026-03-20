from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# MIP-like palette from current UI styles.
MIP_BLUE = RGBColor(48, 105, 240)
MIP_DARK = RGBColor(15, 23, 42)
MIP_SLATE = RGBColor(51, 65, 85)
MIP_BG = RGBColor(245, 248, 252)
MIP_TEXT = RGBColor(31, 41, 55)
MIP_MUTED = RGBColor(95, 107, 122)
MIP_GREEN = RGBColor(46, 125, 50)
MIP_AMBER = RGBColor(146, 95, 0)
MIP_RED = RGBColor(198, 40, 40)
MIP_PURPLE = RGBColor(106, 27, 154)


def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _draw_header(slide, title: str, subtitle: str, section: str = "MIP DEMO") -> None:
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.78))
    top.fill.solid()
    top.fill.fore_color.rgb = MIP_DARK
    top.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.75), Inches(13.333), Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()

    section_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.14), Inches(2.3), Inches(0.42))
    section_box.fill.solid()
    section_box.fill.fore_color.rgb = RGBColor(36, 48, 71)
    section_box.line.color.rgb = MIP_BLUE
    sec_tf = section_box.text_frame
    sec_tf.text = section
    sec_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    sec_tf.paragraphs[0].font.size = Pt(11)
    sec_tf.paragraphs[0].font.bold = True
    sec_tf.paragraphs[0].font.color.rgb = RGBColor(225, 233, 246)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.02), Inches(8.7), Inches(0.7))
    title_tf = title_box.text_frame
    title_tf.text = title
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = MIP_DARK

    subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.72), Inches(11.8), Inches(0.55))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = subtitle
    subtitle_tf.paragraphs[0].font.size = Pt(16)
    subtitle_tf.paragraphs[0].font.color.rgb = MIP_MUTED


def _add_card(slide, x: float, y: float, w: float, h: float, title: str, items: list[str], edge: RGBColor = MIP_BLUE) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(215, 220, 228)
    card.shadow.inherit = False

    edge_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    edge_bar.fill.solid()
    edge_bar.fill.fore_color.rgb = edge
    edge_bar.line.fill.background()

    tf = card.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(15)
    p0.font.color.rgb = MIP_SLATE
    for item in items:
        p = tf.add_paragraph()
        p.text = f"- {item}"
        p.level = 0
        p.font.size = Pt(12)
        p.font.color.rgb = MIP_TEXT


def _add_chip(slide, x: float, y: float, w: float, text: str, bg: RGBColor, fg: RGBColor = RGBColor(255, 255, 255)) -> None:
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    chip.fill.solid()
    chip.fill.fore_color.rgb = bg
    chip.line.fill.background()
    tf = chip.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = fg


def _add_process(slide, labels: list[str], y: float) -> None:
    start_x = 0.75
    box_w = 1.88
    for idx, label in enumerate(labels):
        x = start_x + idx * 2.1
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = RGBColor(207, 214, 225)
        tf = box.text_frame
        tf.text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = MIP_SLATE
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + box_w + 0.06), Inches(y + 0.25), Inches(0.18), Inches(0.45))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MIP_BLUE
            arrow.line.fill.background()


def _add_footer(slide, text: str = "MIP | Evidence -> Decision -> Review") -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.32))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(238, 242, 248)
    bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.6), Inches(7.22), Inches(10.0), Inches(0.2))
    tf = t.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = RGBColor(110, 123, 141)


def _add_visual_title_slide(prs: Presentation, splash_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Market Intelligence Platform",
        "Demo + Training | Research Internals to Trading Operations",
        "MIP OVERVIEW",
    )

    if splash_path.exists():
        slide.shapes.add_picture(str(splash_path), Inches(7.2), Inches(1.45), Inches(5.6), Inches(4.8))

    hero = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.2), Inches(6.2), Inches(3.8))
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(255, 255, 255)
    hero.line.color.rgb = RGBColor(215, 220, 228)
    tf = hero.text_frame
    tf.text = "Why MIP"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MIP_DARK
    for txt in [
        "Evidence-first workflow, not prediction theater.",
        "Committee + risk gates keep execution controlled.",
        "Every decision is auditable from signal to outcome.",
    ]:
        p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(16)
        p.font.color.rgb = MIP_TEXT

    _add_chip(slide, 0.8, 6.2, 2.0, "RESEARCH", MIP_BLUE)
    _add_chip(slide, 2.95, 6.2, 2.0, "DECISIONS", MIP_PURPLE)
    _add_chip(slide, 5.1, 6.2, 2.0, "TRADING", MIP_GREEN)
    _add_footer(slide)


def _add_agenda_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Demo Flow", "Inside-out sequence: research -> committee -> live operations", "AGENDA")

    phases = [
        ("1", "Research Foundations", "Intro, Training, Timeline, Parallel Worlds, News", MIP_BLUE),
        ("2", "Decision Layer", "AI Committee decisions + symbol commentary", MIP_PURPLE),
        ("3", "Trading Operations", "Live link, activity, performance, audit, ledger", MIP_GREEN),
    ]
    x = 0.8
    for num, title, desc, color in phases:
        bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.2), Inches(0.75), Inches(0.75))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = color
        bubble.line.fill.background()
        bt = bubble.text_frame
        bt.text = num
        bt.paragraphs[0].alignment = PP_ALIGN.CENTER
        bt.paragraphs[0].font.size = Pt(20)
        bt.paragraphs[0].font.bold = True
        bt.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        _add_card(slide, x + 0.9, 2.0, 10.9, 1.3, title, [desc], edge=color)
        x += 0.0
        if num != "3":
            connector = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(3.45 if num == "1" else 4.73), Inches(0.5), Inches(0.35))
            connector.fill.solid()
            connector.fill.fore_color.rgb = RGBColor(183, 194, 214)
            connector.line.fill.background()

        # Move next card down
        x = 0.8
        if num == "1":
            phases_y = 3.45
        elif num == "2":
            phases_y = 4.73
        else:
            phases_y = 6.01
        if num != "1":
            bubble.left = Inches(0.8)
            bubble.top = Inches(phases_y)
            slide.shapes._spTree.remove(bubble._element)
            bubble = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(phases_y), Inches(0.75), Inches(0.75))
            bubble.fill.solid()
            bubble.fill.fore_color.rgb = color
            bubble.line.fill.background()
            bt = bubble.text_frame
            bt.text = num
            bt.paragraphs[0].alignment = PP_ALIGN.CENTER
            bt.paragraphs[0].font.size = Pt(20)
            bt.paragraphs[0].font.bold = True
            bt.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.7), Inches(phases_y - 0.2), Inches(10.9), Inches(1.3))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = RGBColor(215, 220, 228)
            edge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.7), Inches(phases_y - 0.2), Inches(0.08), Inches(1.3))
            edge.fill.solid()
            edge.fill.fore_color.rgb = color
            edge.line.fill.background()
            tf = card.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(15)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = MIP_SLATE
            p = tf.add_paragraph()
            p.text = f"- {desc}"
            p.font.size = Pt(12)
            p.font.color.rgb = MIP_TEXT

    _add_footer(slide)


def _add_pipeline_intro(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "MIP Introduction", "Research engine + trading controls in one Snowflake-native system", "INTRO")

    _add_process(
        slide,
        ["Market Data", "Signals", "Outcomes", "Trust", "Proposals", "Execution"],
        y=2.2,
    )
    _add_card(
        slide,
        0.8,
        3.6,
        5.9,
        2.9,
        "Research Functionality",
        [
            "Pattern detection with horizon-based evaluation.",
            "Training maturity, coverage, trust by symbol/pattern.",
            "Parallel worlds for policy stress tests.",
            "News context layered as evidence.",
        ],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        6.7,
        3.6,
        5.9,
        2.9,
        "Trading Functionality",
        [
            "Committee verdicts before execution.",
            "Live link controls + readiness gates.",
            "Lifecycle activity and symbol monitoring.",
            "Performance + audit + learning loop.",
        ],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_training_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Training Status", "Evidence quality before action", "RESEARCH")

    _add_card(
        slide,
        0.8,
        2.1,
        6.3,
        2.1,
        "How To Read It",
        [
            "Start with sample size + coverage.",
            "Then compare Avg H1/H3/H5/H10/H20.",
            "Maturity = confidence in evidence, not certainty.",
        ],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        7.2,
        2.1,
        5.4,
        2.1,
        "Example",
        [
            "AUD/USD | Pattern 2",
            "Coverage 86% | Maturity 78",
            "Avg H5 +0.8% (historical edge)",
        ],
        edge=MIP_GREEN,
    )
    _add_chip(slide, 0.8, 4.5, 2.0, "LOW EVIDENCE", MIP_RED)
    _add_chip(slide, 3.0, 4.5, 2.0, "LEARNING", MIP_AMBER)
    _add_chip(slide, 5.2, 4.5, 2.0, "TRUST READY", MIP_GREEN)
    _add_card(
        slide,
        0.8,
        5.0,
        11.8,
        1.5,
        "AI Fusion",
        ["Ask MIP translates maturity and horizon stats into plain-language confidence guidance."],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_market_timeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Market Timeline", "Symbol-level observability from signal to trade", "RESEARCH")

    _add_process(slide, ["Signals (S:12)", "Proposals (P:3)", "Trades (T:1)"], y=2.15)
    _add_card(
        slide,
        0.8,
        3.45,
        6.0,
        3.0,
        "Visual Storyboard",
        [
            "Tile colors: executed / proposed / signal-only / inactive.",
            "Chart overlays: signal (blue), proposal (orange), trade (green).",
            "Signal chain tree exposes proposal-by-portfolio branches.",
        ],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        6.9,
        3.45,
        5.7,
        3.0,
        "Narrative Value",
        [
            "Decision narrative explains why a symbol advanced or stalled.",
            "Cross-check trust status by pattern.",
            "Link directly into committee decisions and runs.",
        ],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_parallel_worlds_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Parallel Worlds", "Counterfactual lab for policy quality and regret", "RESEARCH")

    health = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(4.1), Inches(2.0))
    health.fill.solid()
    health.fill.fore_color.rgb = RGBColor(230, 245, 234)
    health.line.color.rgb = RGBColor(165, 214, 167)
    tf = health.text_frame
    tf.text = "Policy Health: HEALTHY"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(17)
    tf.paragraphs[0].font.color.rgb = MIP_GREEN
    for row in ["Stability 95/100", "Regret Driver: Baseline", "Top Candidate: Weak (monitor)"]:
        p = tf.add_paragraph()
        p.text = row
        p.font.size = Pt(12)
        p.font.color.rgb = MIP_TEXT

    _add_card(
        slide,
        5.1,
        2.0,
        7.5,
        2.0,
        "Scenario Types",
        ["Signal filter", "Position sizing", "Entry timing", "Baseline (stay cash)"],
        edge=MIP_BLUE,
    )
    _add_chip(slide, 0.8, 4.3, 1.5, "STRONG", MIP_GREEN)
    _add_chip(slide, 2.4, 4.3, 1.5, "EMERGING", MIP_BLUE)
    _add_chip(slide, 4.0, 4.3, 1.5, "WEAK", MIP_AMBER)
    _add_chip(slide, 5.6, 4.3, 1.5, "NOISE", RGBColor(113, 128, 150))
    _add_card(
        slide,
        0.8,
        4.85,
        11.8,
        1.65,
        "AI Fusion",
        ["AI narrative explains gate-by-gate divergence and regret trend; humans still approve all policy changes."],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_news_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "News Intelligence", "Evidence-backed context layer with strict guardrails", "RESEARCH")

    _add_card(
        slide,
        0.8,
        2.0,
        3.7,
        2.8,
        "Context KPIs",
        ["Symbols with news", "HOT symbols", "Snapshot age", "Stale symbols"],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        4.7,
        2.0,
        4.0,
        2.8,
        "Decision Impact",
        ["Proposals scoped", "With news context", "With score adjustment", "Blocked new entry"],
        edge=MIP_PURPLE,
    )
    _add_card(
        slide,
        8.9,
        2.0,
        3.7,
        2.8,
        "Guardrails",
        ["Invalid URLs excluded", "Freshness explicit", "Influence bounded", "No narrative guesswork"],
        edge=MIP_GREEN,
    )
    _add_chip(slide, 0.8, 5.1, 1.2, "HOT", MIP_RED)
    _add_chip(slide, 2.1, 5.1, 2.8, "UNSEEN + DECISION RELEVANT", MIP_DARK)
    _add_card(
        slide,
        0.8,
        5.55,
        11.8,
        0.95,
        "AI Fusion",
        ["Reader summary is generated from stored features and linked proposal payload evidence."],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_live_link_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Live Portfolio Link", "Control plane binding MIP workflow to broker truth", "TRADING")

    _add_process(slide, ["Source Portfolio", "MIP Live Portfolio", "IBKR Account"], y=2.2)
    _add_chip(slide, 4.95, 3.35, 1.8, "Activation Guard", MIP_AMBER)
    _add_chip(slide, 6.9, 3.35, 2.1, "Execution Readiness", MIP_GREEN)

    _add_card(
        slide,
        0.8,
        4.0,
        5.7,
        2.4,
        "Configuration Scope",
        ["Adapter mode", "Position and exposure limits", "Freshness + cooldown controls", "Drawdown and bust brakes"],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        6.9,
        4.0,
        5.7,
        2.4,
        "Governance",
        ["Saving writes control state only.", "No order placement from this page.", "Execution still requires validation and approvals."],
        edge=MIP_GREEN,
    )
    _add_footer(slide)


def _add_live_activity_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Live Portfolio Activity", "Operational lifecycle with committee checkpoints", "TRADING")

    _add_process(slide, ["Imported", "Validated", "Committee", "Approved", "Executed"], y=2.2)
    _add_chip(slide, 5.0, 3.35, 2.3, "Committee checkpoint", MIP_PURPLE)
    _add_card(
        slide,
        0.8,
        4.0,
        5.8,
        2.4,
        "What You Verify",
        ["Recent transitions", "State freshness", "Delay/block reasons", "Approval path integrity"],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        6.8,
        4.0,
        5.8,
        2.4,
        "Cross-Checks",
        ["AI Agent Decisions for verdict context", "Runs for execution truth", "Symbol Tracker for live thesis state"],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_decisions_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "AI Agent Decisions", "Committee courtroom for verdict, reason codes, and revalidation", "DECISIONS")

    _add_card(
        slide,
        0.8,
        2.1,
        5.8,
        2.1,
        "Accepted Example",
        ["Verdict: APPROVE", "Confidence: HIGH", "Reason: Trust + risk + freshness all passed"],
        edge=MIP_GREEN,
    )
    _add_card(
        slide,
        6.8,
        2.1,
        5.8,
        2.1,
        "Rejected Example",
        ["Verdict: REJECT", "Confidence: MEDIUM", "Reason: Capacity/risk gate block, stale quote risk"],
        edge=MIP_RED,
    )
    _add_process(slide, ["PROPOSED", "APPROVED", "REJECTED / EXECUTED"], y=4.6)
    _add_card(
        slide,
        0.8,
        5.55,
        11.8,
        0.95,
        "AI Fusion",
        ["Committee summary explains verdict quickly; reason codes preserve deterministic accountability."],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_symbol_tracker_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, RGBColor(11, 18, 32))
    _draw_header(slide, "Live Symbol Tracker", "Symbol-first monitoring + committee commentary", "DECISIONS")

    _add_chip(slide, 0.9, 2.2, 2.2, "THESIS_INTACT", RGBColor(3, 105, 70))
    _add_chip(slide, 3.3, 2.2, 2.0, "WEAKENING", MIP_AMBER)
    _add_chip(slide, 5.5, 2.2, 2.1, "INVALIDATED", MIP_RED)

    dark_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.8), Inches(5.8), Inches(3.7))
    dark_card.fill.solid()
    dark_card.fill.fore_color.rgb = RGBColor(15, 23, 42)
    dark_card.line.color.rgb = RGBColor(36, 48, 65)
    tf = dark_card.text_frame
    tf.text = "Symbol Metrics"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
    tf.paragraphs[0].font.size = Pt(16)
    for row in ["Open R: +1.2R", "Expected move reached: 84%", "Distance to SL: 1.6%", "Vol regime: Normal -> Elevated"]:
        p = tf.add_paragraph()
        p.text = f"- {row}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(203, 213, 225)

    dark_card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.8), Inches(5.8), Inches(3.7))
    dark_card2.fill.solid()
    dark_card2.fill.fore_color.rgb = RGBColor(15, 23, 42)
    dark_card2.line.color.rgb = RGBColor(36, 48, 65)
    tf2 = dark_card2.text_frame
    tf2.text = "Committee Commentary"
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
    tf2.paragraphs[0].font.size = Pt(16)
    for row in [
        "Stance: WATCH_CLOSELY",
        "Confidence: MEDIUM",
        "Reason tags: momentum_fade, near_tp, news_uncertainty",
        "Actions to consider: tighten stop, monitor revalidation",
    ]:
        p = tf2.add_paragraph()
        p.text = f"- {row}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(203, 213, 225)

    _add_footer(slide, text="MIP | Symbol monitoring + committee commentary")


def _add_learning_ledger_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Learning Ledger", "Causality trace: what changed, why, and what happened next", "REVIEW")

    _add_process(slide, ["Proposal", "Decision", "Execution", "Outcome", "Lesson"], y=2.2)
    _add_card(
        slide,
        0.8,
        3.65,
        5.8,
        2.8,
        "What It Proves",
        ["Evidence chain behind each decision", "Role/context attribution", "Expected vs realized outcome", "Recurring pattern quality over time"],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        6.8,
        3.65,
        5.8,
        2.8,
        "How Teams Use It",
        ["Post-trade review", "Weekly retrospective", "Policy tuning and rejection pattern analysis", "Stakeholder explainability"],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_performance_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Performance Dashboard", "Portfolio-level truth on return, drawdown, and consistency", "REVIEW")

    _add_card(
        slide,
        0.8,
        2.0,
        3.8,
        2.5,
        "Portfolio A",
        ["Return: +12.4%", "Max DD: -4.1%", "Stability: High"],
        edge=MIP_GREEN,
    )
    _add_card(
        slide,
        4.8,
        2.0,
        3.8,
        2.5,
        "Portfolio B",
        ["Return: +8.7%", "Max DD: -2.9%", "Stability: Very High"],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        8.8,
        2.0,
        3.8,
        2.5,
        "Portfolio C",
        ["Return: +15.1%", "Max DD: -7.3%", "Stability: Medium"],
        edge=MIP_AMBER,
    )
    _add_card(
        slide,
        0.8,
        4.8,
        11.8,
        1.7,
        "Interpretation Workflow",
        ["Use period filters first, then compare return vs drawdown, then trace root cause in Training + Decisions + Runs."],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_audit_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "Audit Trail (Runs)", "Operational truth: what ran, status, durations, and errors", "OPERATIONS")

    _add_chip(slide, 0.8, 2.05, 1.3, "SUCCESS", MIP_GREEN)
    _add_chip(slide, 2.2, 2.05, 1.7, "SUCCESS+SKIPS", MIP_AMBER)
    _add_chip(slide, 4.0, 2.05, 1.2, "FAILED", MIP_RED)
    _add_chip(slide, 5.3, 2.05, 1.4, "RUNNING", MIP_BLUE)

    _add_card(
        slide,
        0.8,
        2.6,
        5.8,
        3.9,
        "Run Detail Panel",
        ["Summary cards: status, duration, as-of, portfolios.", "Step timeline with per-step duration and outcome.", "Error panel with SQLSTATE and query IDs."],
        edge=MIP_BLUE,
    )
    _add_card(
        slide,
        6.8,
        2.6,
        5.8,
        3.9,
        "Demo Note",
        ["Open one successful run and one failed run.", "Show how AI summary accelerates triage.", "Confirm final truth in step-level diagnostics."],
        edge=MIP_PURPLE,
    )
    _add_footer(slide)


def _add_ai_fusion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(slide, "AI Is Naturally Fused", "AI accelerates understanding while controls keep humans in charge", "GOVERNANCE")

    _add_card(
        slide,
        0.8,
        2.0,
        5.8,
        4.7,
        "Where AI Helps",
        ["Daily digests and committee summaries.", "Parallel Worlds what-if explanation.", "News reader summary and HOT prioritization.", "Ask MIP route-aware training support."],
        edge=MIP_PURPLE,
    )
    _add_card(
        slide,
        6.8,
        2.0,
        5.8,
        4.7,
        "Where Humans Control",
        ["Risk gates are deterministic.", "Execution needs readiness + approvals.", "AI comments are non-binding.", "Runs + Ledger preserve accountability."],
        edge=MIP_GREEN,
    )
    _add_footer(slide)


def _add_close_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_DARK)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = MIP_DARK
    banner.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.3), Inches(0.15), Inches(4.8))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.3), Inches(1.5), Inches(10.5), Inches(1.4))
    tf = title.text_frame
    tf.text = "Q&A and Live Walkthrough"
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(241, 245, 249)

    sub = slide.shapes.add_textbox(Inches(1.3), Inches(3.1), Inches(10.8), Inches(2.0))
    sf = sub.text_frame
    sf.text = "Suggested final flow in front of colleagues:"
    sf.paragraphs[0].font.size = Pt(20)
    sf.paragraphs[0].font.color.rgb = RGBColor(203, 213, 225)
    for row in [
        "1) Training -> Timeline -> Parallel Worlds",
        "2) News -> AI Agent Decisions -> Symbol Tracker commentary",
        "3) Live Link/Activity -> Performance -> Runs/Audit -> Learning Ledger",
    ]:
        p = sf.add_paragraph()
        p.text = row
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(226, 232, 240)


def build_deck(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    splash = Path(__file__).resolve().parents[1] / "apps" / "mip_ui_web" / "public" / "mip-splash.png"

    _add_visual_title_slide(prs, splash)
    _add_agenda_slide(prs)
    _add_pipeline_intro(prs)
    _add_training_slide(prs)
    _add_market_timeline_slide(prs)
    _add_parallel_worlds_slide(prs)
    _add_news_slide(prs)
    _add_live_link_slide(prs)
    _add_live_activity_slide(prs)
    _add_decisions_slide(prs)
    _add_symbol_tracker_slide(prs)
    _add_learning_ledger_slide(prs)
    _add_performance_slide(prs)
    _add_audit_slide(prs)
    _add_ai_fusion_slide(prs)
    _add_close_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    deck_path = Path(__file__).with_name("MIP_Demo_Training_Deck_Styled.pptx")
    build_deck(deck_path)
    print(f"Created: {deck_path}")
