"""
Build MIP Platform Introduction Presentation
Market Research Platform + Broker-Linked Trading

Uses the MIP design system: dark headers, blue accents, white cards, edge bars,
process flows, chips, and visual analogies throughout.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── MIP palette ──────────────────────────────────────────────────────────────
MIP_BLUE      = RGBColor(48, 105, 240)
MIP_DARK      = RGBColor(15, 23, 42)
MIP_SLATE     = RGBColor(51, 65, 85)
MIP_BG        = RGBColor(245, 248, 252)
MIP_TEXT      = RGBColor(31, 41, 55)
MIP_MUTED     = RGBColor(95, 107, 122)
MIP_GREEN     = RGBColor(46, 125, 50)
MIP_AMBER     = RGBColor(146, 95, 0)
MIP_RED       = RGBColor(198, 40, 40)
MIP_PURPLE    = RGBColor(106, 27, 154)
MIP_TEAL      = RGBColor(0, 121, 107)
MIP_CYAN      = RGBColor(0, 151, 167)
MIP_INDIGO    = RGBColor(57, 73, 171)
WHITE         = RGBColor(255, 255, 255)
LIGHT_GRAY    = RGBColor(215, 220, 228)
CARD_BG       = RGBColor(255, 255, 255)
DARK_CARD_BG  = RGBColor(15, 23, 42)
DARK_CARD_BDR = RGBColor(36, 48, 65)
HERO_GREEN_BG = RGBColor(230, 245, 234)
HERO_BLUE_BG  = RGBColor(227, 236, 253)
HERO_PURP_BG  = RGBColor(243, 229, 245)
HERO_AMB_BG   = RGBColor(255, 248, 225)
FOOTER_BG     = RGBColor(238, 242, 248)
FOOTER_TEXT    = RGBColor(110, 123, 141)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper functions ─────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _draw_header(slide, title: str, subtitle: str, section: str = "MIP PLATFORM") -> None:
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.78))
    top.fill.solid()
    top.fill.fore_color.rgb = MIP_DARK
    top.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.75), SLIDE_W, Inches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()

    sec_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.14), Inches(2.3), Inches(0.42)
    )
    sec_box.fill.solid()
    sec_box.fill.fore_color.rgb = RGBColor(36, 48, 71)
    sec_box.line.color.rgb = MIP_BLUE
    tf = sec_box.text_frame
    tf.text = section
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(225, 233, 246)

    t_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.02), Inches(9.5), Inches(0.7))
    tf = t_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(34)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MIP_DARK

    s_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.72), Inches(11.8), Inches(0.55))
    tf = s_box.text_frame
    tf.text = subtitle
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = MIP_MUTED


def _add_card(
    slide, x: float, y: float, w: float, h: float,
    title: str, items: list[str], edge: RGBColor = MIP_BLUE,
    title_size: int = 15, item_size: int = 12,
) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = LIGHT_GRAY
    card.shadow.inherit = False

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = edge
    bar.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(title_size)
    p0.font.color.rgb = MIP_SLATE
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(item_size)
        p.font.color.rgb = MIP_TEXT
        p.space_before = Pt(4)


def _add_dark_card(
    slide, x: float, y: float, w: float, h: float,
    title: str, items: list[str], edge: RGBColor = MIP_BLUE,
) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_CARD_BG
    card.line.color.rgb = DARK_CARD_BDR

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = edge
    bar.line.fill.background()

    tf = card.text_frame
    tf.word_wrap = True
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.bold = True
    p0.font.size = Pt(15)
    p0.font.color.rgb = RGBColor(226, 232, 240)
    for item in items:
        p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_before = Pt(3)


def _add_chip(
    slide, x: float, y: float, w: float, text: str,
    bg: RGBColor, fg: RGBColor = WHITE, h: float = 0.34,
) -> None:
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    chip.fill.solid()
    chip.fill.fore_color.rgb = bg
    chip.line.fill.background()
    tf = chip.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = fg


def _add_process(slide, labels: list[str], y: float, start_x: float = 0.75, box_w: float = 1.88, gap: float = 0.22) -> None:
    for idx, label in enumerate(labels):
        x = start_x + idx * (box_w + gap)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(0.85)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = RGBColor(207, 214, 225)
        tf = box.text_frame
        tf.text = label
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = MIP_SLATE
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(x + box_w + 0.02), Inches(y + 0.2),
                Inches(0.18), Inches(0.45),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MIP_BLUE
            arrow.line.fill.background()


def _add_icon_circle(
    slide, x: float, y: float, size: float,
    text: str, bg: RGBColor, fg: RGBColor = WHITE,
    font_size: int = 22,
) -> None:
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = bg
    circle.line.fill.background()
    tf = circle.text_frame
    tf.text = text
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = fg


def _add_hero_box(
    slide, x: float, y: float, w: float, h: float,
    bg: RGBColor, title: str, body_lines: list[str],
    title_color: RGBColor = MIP_DARK,
) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg
    box.line.color.rgb = LIGHT_GRAY
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = title
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = title_color
    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = MIP_TEXT
        p.space_before = Pt(4)


def _add_analogy_bar(slide, icon: str, text: str, y: float = 6.6) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(240, 243, 250)
    bar.line.color.rgb = LIGHT_GRAY
    tf = bar.text_frame
    tf.text = f"{icon}  {text}"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = MIP_SLATE


def _add_footer(slide, text: str = "MIP  |  Market Intelligence Platform  |  Evidence  >  Decision  >  Execution") -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), SLIDE_W, Inches(0.32)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOOTER_BG
    bar.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.6), Inches(7.22), Inches(10.0), Inches(0.2))
    tf = t.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = FOOTER_TEXT


def _add_stat_box(
    slide, x: float, y: float, w: float, h: float,
    value: str, label: str, color: RGBColor,
) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = LIGHT_GRAY
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = value
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = color
    p = tf.add_paragraph()
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.color.rgb = MIP_MUTED
    p.space_before = Pt(2)


def _add_gradient_banner(slide, x, y, w, h, color1, color2):
    """Simulate gradient with two overlapping rectangles."""
    r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w/2), Inches(h))
    r1.fill.solid()
    r1.fill.fore_color.rgb = color1
    r1.line.fill.background()
    r2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + w/2), Inches(y), Inches(w/2), Inches(h))
    r2.fill.solid()
    r2.fill.fore_color.rgb = color2
    r2.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def _slide_title(prs: Presentation) -> None:
    """Slide 1: Grand opening — What is MIP."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_DARK)

    _add_gradient_banner(slide, 0, 0, 13.333, 7.5, MIP_DARK, RGBColor(20, 30, 58))

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.0), Inches(0.15), Inches(5.2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.3), Inches(1.2), Inches(7.5), Inches(1.0))
    tf = title.text_frame
    tf.text = "Market Intelligence Platform"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    sub1 = slide.shapes.add_textbox(Inches(1.3), Inches(2.3), Inches(8.0), Inches(0.6))
    tf = sub1.text_frame
    tf.text = "Where Market Research Meets Intelligent Execution"
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(3.1), Inches(3.0), Inches(0.04)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = MIP_BLUE
    divider.line.fill.background()

    desc = slide.shapes.add_textbox(Inches(1.3), Inches(3.4), Inches(6.5), Inches(2.5))
    tf = desc.text_frame
    tf.word_wrap = True
    lines = [
        "MIP is a Snowflake-native platform that fuses AI-driven",
        "market research with live broker-linked trading.",
        "",
        "Every trade is backed by evidence. Every decision is auditable.",
        "Every outcome feeds the next cycle of learning.",
    ]
    tf.text = lines[0]
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.color.rgb = RGBColor(203, 213, 225)
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(203, 213, 225)
        p.space_before = Pt(4)

    pillars = [
        ("RESEARCH", MIP_BLUE),
        ("DECISIONS", MIP_PURPLE),
        ("EXECUTION", MIP_GREEN),
        ("LEARNING", MIP_TEAL),
    ]
    for i, (label, color) in enumerate(pillars):
        _add_chip(slide, 1.3 + i * 2.2, 6.1, 1.9, label, color)

    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), Inches(1.5), Inches(4.5), Inches(4.8)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = RGBColor(20, 32, 55)
    right_card.line.color.rgb = RGBColor(40, 58, 90)
    tf = right_card.text_frame
    tf.word_wrap = True
    tf.text = "Platform Highlights"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)
    highlights = [
        "AI committee reviews every trade proposal",
        "Pattern-based signal detection across FX + equities",
        "Parallel worlds test policy quality before risk",
        "Live IBKR link with real-time position sync",
        "News intelligence with bounded influence",
        "Full audit trail from signal to settlement",
        "Learning ledger closes the feedback loop",
    ]
    for h in highlights:
        p = tf.add_paragraph()
        p.text = f"  {h}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(180, 195, 215)
        p.space_before = Pt(6)


def _slide_what_is_mip(prs: Presentation) -> None:
    """Slide 2: What Is MIP — the elevator pitch with visual pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "What Is MIP?",
        "An evidence-first market intelligence engine — not a prediction machine",
        "OVERVIEW",
    )

    _add_process(
        slide,
        ["Market Data", "Signals", "Training", "Proposals", "Committee", "Execution"],
        y=2.15, start_x=0.5, box_w=1.82, gap=0.25,
    )

    _add_hero_box(
        slide, 0.7, 3.4, 5.9, 3.0, HERO_BLUE_BG,
        "Research Engine",
        [
            "Ingests market data across FX pairs and equities",
            "Detects patterns and evaluates across 5 horizons (H1-H20)",
            "Builds trust scores through backtested evidence",
            "News context + parallel world stress testing",
            "Training maturity gates prevent premature trading",
        ],
        title_color=MIP_BLUE,
    )

    _add_hero_box(
        slide, 6.8, 3.4, 5.9, 3.0, HERO_GREEN_BG,
        "Trading Operations",
        [
            "AI committee verdicts before any execution",
            "Live broker link with IBKR (paper + live modes)",
            "Position lifecycle: import > validate > approve > execute",
            "Performance dashboards with Sharpe, drawdown, win rate",
            "Every decision traceable in the Learning Ledger",
        ],
        title_color=MIP_GREEN,
    )

    _add_analogy_bar(
        slide, "\u2696\uFE0F",
        "Think of MIP as a research lab with a built-in trading desk — the lab must approve before the desk can act.",
    )
    _add_footer(slide)


def _slide_training_status(prs: Presentation) -> None:
    """Slide 3: Training Status."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Training Status",
        "Evidence quality is measured before any trade is considered — maturity gates keep the system honest",
        "RESEARCH",
    )

    _add_stat_box(slide, 0.7, 2.15, 2.8, 1.3, "0-100", "Maturity Score", MIP_BLUE)
    _add_stat_box(slide, 3.6, 2.15, 2.8, 1.3, "H1-H20", "5 Horizons", MIP_PURPLE)
    _add_stat_box(slide, 6.5, 2.15, 2.8, 1.3, "86%", "Coverage Target", MIP_GREEN)
    _add_stat_box(slide, 9.4, 2.15, 3.2, 1.3, "PER SYMBOL", "Granularity", MIP_TEAL)

    _add_chip(slide, 0.7, 3.7, 2.3, "INSUFFICIENT", MIP_RED)
    _add_chip(slide, 3.1, 3.7, 2.3, "WARMING UP", MIP_AMBER)
    _add_chip(slide, 5.5, 3.7, 2.3, "LEARNING", MIP_BLUE)
    _add_chip(slide, 7.9, 3.7, 2.3, "CONFIDENT", MIP_GREEN)

    _add_card(
        slide, 0.7, 4.3, 6.0, 2.2,
        "How Training Works",
        [
            "Each symbol + pattern combination is independently trained",
            "Maturity = sample size + coverage + horizon depth",
            "Trust labels gate whether the system can propose trades",
            "LOW_EVIDENCE symbols are quarantined from execution",
            "The system learns continuously — maturity rises over time",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 6.85, 4.3, 5.75, 2.2,
        "What You See in the UI",
        [
            "Symbol grid with maturity scores and trust labels",
            "Horizon returns: Avg H1, H3, H5, H10, H20 per pattern",
            "Coverage % — how much of the symbol's history is trained",
            "Sample size — number of observations backing the score",
            "AI summary translates stats into plain-language guidance",
        ],
        edge=MIP_PURPLE,
    )

    _add_analogy_bar(
        slide, "\U0001F393",
        "Like a medical residency: doctors don't operate until they've logged enough supervised hours. MIP's patterns don't trade until trained.",
    )
    _add_footer(slide)


def _slide_market_timeline(prs: Presentation) -> None:
    """Slide 4: Market Timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Market Timeline",
        "Symbol-level storyboard — see every signal, proposal, and trade on a visual timeline",
        "RESEARCH",
    )

    _add_process(
        slide,
        ["Signals (S:12)", "Proposals (P:3)", "Trades (T:1)"],
        y=2.15, start_x=1.5, box_w=3.0, gap=0.5,
    )

    _add_chip(slide, 0.7, 3.25, 1.6, "EXECUTED", MIP_GREEN)
    _add_chip(slide, 2.4, 3.25, 1.6, "PROPOSED", MIP_BLUE)
    _add_chip(slide, 4.1, 3.25, 1.8, "SIGNAL ONLY", MIP_AMBER)
    _add_chip(slide, 6.0, 3.25, 1.6, "INACTIVE", RGBColor(113, 128, 150))

    _add_card(
        slide, 0.7, 3.85, 6.0, 2.7,
        "Visual Storyboard",
        [
            "Tile grid colors indicate lifecycle stage at a glance",
            "Chart overlays: signals (blue), proposals (orange), trades (green)",
            "Signal chain tree shows branching proposals per portfolio",
            "OHLC price chart with signal/trade markers overlaid",
            "Cross-check training trust status by pattern directly",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 6.85, 3.85, 5.75, 2.7,
        "Narrative Value",
        [
            "Decision narrative explains why a symbol advanced or stalled",
            "Full transparency: signal > proposal > committee > outcome",
            "Click-through to committee decisions and execution runs",
            "Portfolio-level filtering for multi-strategy environments",
            "Time-range controls for historical investigation",
        ],
        edge=MIP_PURPLE,
    )

    _add_analogy_bar(
        slide, "\U0001F4F0",
        "Like a newspaper front page for each symbol — you see the headline (trade), the story (proposals), and the raw sources (signals).",
    )
    _add_footer(slide)


def _slide_news_intelligence(prs: Presentation) -> None:
    """Slide 5: News Intelligence."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "News Intelligence",
        "Evidence-backed context layer — news informs decisions but never overrides risk gates",
        "RESEARCH",
    )

    _add_stat_box(slide, 0.7, 2.15, 2.85, 1.3, "RSS", "Multi-Source Ingest", MIP_BLUE)
    _add_stat_box(slide, 3.65, 2.15, 2.85, 1.3, "Z \u2265 1.5", "HOT Threshold", MIP_RED)
    _add_stat_box(slide, 6.6, 2.15, 2.85, 1.3, "BOUNDED", "Influence Control", MIP_PURPLE)
    _add_stat_box(slide, 9.55, 2.15, 3.1, 1.3, "REAL-TIME", "Freshness Tracking", MIP_GREEN)

    _add_card(
        slide, 0.7, 3.7, 3.9, 2.8,
        "Context KPIs",
        [
            "Symbols with active news coverage",
            "HOT symbols (burst z-score \u2265 1.5)",
            "Snapshot age and freshness tracking",
            "Stale symbol alerts for decision caution",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 4.75, 3.7, 3.9, 2.8,
        "Decision Impact",
        [
            "Proposals scoped with news context",
            "Score adjustments bounded by guardrails",
            "New entry blocks when uncertainty spikes",
            "AI reader summary from stored features",
        ],
        edge=MIP_PURPLE,
    )

    _add_card(
        slide, 8.8, 3.7, 3.85, 2.8,
        "Sources & Guardrails",
        [
            "SEC, GlobeNewswire, MarketWatch, Fed, ECB",
            "IBKR research feed integration",
            "Invalid URLs auto-excluded",
            "No narrative guesswork — evidence only",
        ],
        edge=MIP_GREEN,
    )

    _add_analogy_bar(
        slide, "\U0001F4E1",
        "Like a newsroom fact-checker: MIP reads the headlines but only acts on verified, scored evidence — never on rumor alone.",
    )
    _add_footer(slide)


def _slide_parallel_worlds(prs: Presentation) -> None:
    """Slide 6: Parallel Worlds."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Parallel Worlds",
        "Counterfactual laboratory — stress-test policy quality against 'what if' scenarios before risking capital",
        "RESEARCH",
    )

    health = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.1), Inches(4.3), Inches(2.0)
    )
    health.fill.solid()
    health.fill.fore_color.rgb = HERO_GREEN_BG
    health.line.color.rgb = RGBColor(165, 214, 167)
    tf = health.text_frame
    tf.word_wrap = True
    tf.text = "Policy Health: HEALTHY"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = MIP_GREEN
    for row in [
        "Stability Score: 95 / 100",
        "Regret Driver: Baseline (staying in cash)",
        "Top Candidate: Weak — continue monitoring",
        "Current policy outperforms all alternatives",
    ]:
        p = tf.add_paragraph()
        p.text = f"  {row}"
        p.font.size = Pt(12)
        p.font.color.rgb = MIP_TEXT
        p.space_before = Pt(3)

    _add_card(
        slide, 5.15, 2.1, 7.5, 2.0,
        "Scenario Types Under Test",
        [
            "THRESHOLD — vary signal filter strictness",
            "SIZING — test alternative position sizing models",
            "TIMING — simulate delayed entry by N bars",
            "HORIZON — change holding period assumptions",
            "EARLY_EXIT — apply payoff multiplier cutoffs",
            "BASELINE — the 'do nothing' cash benchmark",
        ],
        edge=MIP_BLUE,
    )

    _add_chip(slide, 0.7, 4.35, 1.5, "STRONG", MIP_GREEN)
    _add_chip(slide, 2.3, 4.35, 1.7, "EMERGING", MIP_BLUE)
    _add_chip(slide, 4.1, 4.35, 1.3, "WEAK", MIP_AMBER)
    _add_chip(slide, 5.5, 4.35, 1.3, "NOISE", RGBColor(113, 128, 150))

    _add_card(
        slide, 0.7, 4.95, 5.9, 1.55,
        "How It Works",
        [
            "Scenarios are simulated over historical data windows",
            "Gate-by-gate divergence analysis identifies where policies differ",
            "Regret = excess return of best alternative vs current policy",
            "Humans approve all policy changes — AI explains, never decides",
        ],
        edge=MIP_TEAL,
    )

    _add_card(
        slide, 6.75, 4.95, 5.9, 1.55,
        "AI Narrative",
        [
            "AI explains gate-by-gate divergence and regret trend",
            "Generates plain-language summary of scenario performance",
            "Highlights emerging threats to current policy stability",
            "Links to live broker snapshots for real-world validation",
        ],
        edge=MIP_PURPLE,
    )

    _add_analogy_bar(
        slide, "\U0001F30D",
        "Like a flight simulator for your trading strategy — test every 'what-if' scenario before you take off with real capital.",
    )
    _add_footer(slide)


def _slide_live_link(prs: Presentation) -> None:
    """Slide 7: Live Portfolio Link."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Live Portfolio Link",
        "The control plane — binding MIP's intelligence engine to real broker execution via Interactive Brokers",
        "TRADING",
    )

    _add_process(
        slide,
        ["Source Portfolio", "MIP Live Config", "IBKR Account"],
        y=2.15, start_x=1.2, box_w=3.2, gap=0.6,
    )

    _add_chip(slide, 4.4, 3.2, 2.0, "Activation Guard", MIP_AMBER)
    _add_chip(slide, 6.6, 3.2, 2.3, "Execution Readiness", MIP_GREEN)

    _add_card(
        slide, 0.7, 3.8, 4.0, 2.7,
        "Configuration",
        [
            "Adapter mode: PAPER or LIVE",
            "Max positions and exposure limits",
            "Cash buffer percentage",
            "Quote freshness thresholds",
            "Cooldown bars between trades",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 4.85, 3.8, 4.0, 2.7,
        "Risk Brakes",
        [
            "Drawdown stop threshold",
            "Bust percentage limit",
            "Validity window enforcement",
            "Snapshot freshness gates",
            "Position-level safeguards",
        ],
        edge=MIP_RED,
    )

    _add_card(
        slide, 9.0, 3.8, 3.65, 2.7,
        "Governance",
        [
            "Save writes config — no orders",
            "Execution requires validation",
            "Committee must approve first",
            "Full audit of every state change",
            "Human override at every step",
        ],
        edge=MIP_GREEN,
    )

    _add_analogy_bar(
        slide, "\U0001F3DB\uFE0F",
        "Like air traffic control: MIP clears each flight (trade) only when conditions, fuel (capital), and weather (market) all pass safety checks.",
    )
    _add_footer(slide)


def _slide_live_activity(prs: Presentation) -> None:
    """Slide 8: Live Portfolio Activity."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Live Portfolio Activity",
        "Real-time operational lifecycle with committee checkpoints at every stage",
        "TRADING",
    )

    _add_process(
        slide,
        ["Imported", "Validated", "Committee", "Approved", "Executed"],
        y=2.15, start_x=0.65, box_w=2.15, gap=0.32,
    )

    _add_chip(slide, 4.95, 3.2, 2.8, "Committee Checkpoint", MIP_PURPLE)

    _add_stat_box(slide, 0.7, 3.75, 3.0, 1.2, "NAV", "Net Asset Value", MIP_GREEN)
    _add_stat_box(slide, 3.85, 3.75, 3.0, 1.2, "P&L", "Unrealized Gains", MIP_BLUE)
    _add_stat_box(slide, 7.0, 3.75, 3.0, 1.2, "DRIFT", "Alignment Check", MIP_AMBER)
    _add_stat_box(slide, 10.15, 3.75, 2.5, 1.2, "W / L", "Winners vs Losers", MIP_PURPLE)

    _add_card(
        slide, 0.7, 5.2, 6.0, 1.4,
        "What You Monitor",
        [
            "Open positions: symbol, side, qty, avg cost, unrealized P&L",
            "Execution history: fill prices, realized P&L, timestamps",
            "TP/SL status + protection level for each position",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 6.85, 5.2, 5.75, 1.4,
        "Cross-Checks Available",
        [
            "AI Agent Decisions  for verdict context and reason codes",
            "Runs  for execution truth and pipeline health",
            "Symbol Tracker  for real-time thesis validation",
        ],
        edge=MIP_PURPLE,
    )

    _add_analogy_bar(
        slide, "\U0001F3E5",
        "Like a hospital patient monitor — vital signs (NAV, P&L, drift) are tracked continuously, and alerts fire before problems escalate.",
    )
    _add_footer(slide)


def _slide_ai_decisions(prs: Presentation) -> None:
    """Slide 9: AI Agent Decisions."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "AI Agent Decisions",
        "The committee courtroom — every trade proposal faces a structured verdict with reason codes",
        "DECISIONS",
    )

    _add_card(
        slide, 0.7, 2.15, 5.9, 2.3,
        "APPROVED Example",
        [
            "Verdict:  APPROVE",
            "Confidence:  HIGH",
            "Trust gate:  PASSED — maturity 82, coverage 91%",
            "Risk gate:  PASSED — within position limits",
            "Freshness:  PASSED — quote age 3 min",
            "Committee summary generated for audit trail",
        ],
        edge=MIP_GREEN,
    )

    _add_card(
        slide, 6.75, 2.15, 5.9, 2.3,
        "REJECTED Example",
        [
            "Verdict:  REJECT",
            "Confidence:  MEDIUM",
            "Trust gate:  FAILED — maturity only 34",
            "Capacity gate:  BLOCKED — max positions reached",
            "Stale quote risk flagged (15 min age)",
            "Reason codes preserved for learning ledger",
        ],
        edge=MIP_RED,
    )

    _add_process(
        slide,
        ["PROPOSED", "TRUST GATE", "RISK GATE", "COMMITTEE", "VERDICT"],
        y=4.75, start_x=0.65, box_w=2.15, gap=0.32,
    )

    _add_card(
        slide, 0.7, 5.85, 11.9, 0.7,
        "Deterministic Accountability",
        [
            "Reason codes are structured and queryable — not free-text opinions. Every APPROVE and REJECT is reproducible and auditable.",
        ],
        edge=MIP_PURPLE, item_size=13,
    )

    _add_analogy_bar(
        slide, "\u2696\uFE0F",
        "Like a courtroom trial: evidence is presented, gates are checked, and a structured verdict is rendered — with a full transcript for appeal.",
    )
    _add_footer(slide)


def _slide_learning_ledger(prs: Presentation) -> None:
    """Slide 10: Learning Ledger."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Learning Ledger",
        "The institutional memory — every decision, its context, and what happened next",
        "REVIEW",
    )

    _add_process(
        slide,
        ["Proposal", "Decision", "Execution", "Outcome", "Lesson"],
        y=2.15, start_x=0.65, box_w=2.15, gap=0.32,
    )

    _add_card(
        slide, 0.7, 3.35, 4.0, 3.1,
        "Event Types",
        [
            "TRAINING_EVENT",
            "  Maturity changes, trust label shifts",
            "DECISION_EVENT",
            "  Committee verdicts, reason codes",
            "LIVE_EVENT",
            "  Execution fills, P&L outcomes",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 4.85, 3.35, 3.9, 3.1,
        "Causality Chain",
        [
            "Before State  vs  After State",
            "Influence delta tracking",
            "Causality links between events",
            "Outcome state for closed-loop review",
            "News context influence measured",
        ],
        edge=MIP_TEAL,
    )

    _add_card(
        slide, 8.9, 3.35, 3.75, 3.1,
        "Team Use Cases",
        [
            "Post-trade retrospectives",
            "Weekly performance reviews",
            "Policy tuning decisions",
            "Rejection pattern analysis",
            "Stakeholder explainability",
        ],
        edge=MIP_PURPLE,
    )

    _add_analogy_bar(
        slide, "\U0001F4DA",
        "Like a pilot's flight log: every takeoff, turbulence, and landing is recorded — so the next flight benefits from every past experience.",
    )
    _add_footer(slide)


def _slide_symbol_tracker(prs: Presentation) -> None:
    """Slide 11: Live Symbol Tracker — dark theme."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, RGBColor(11, 18, 32))
    _draw_header(
        slide,
        "Live Symbol Tracker",
        "Real-time thesis monitoring with committee commentary on every open position",
        "MONITORING",
    )

    _add_chip(slide, 0.7, 2.15, 2.2, "THESIS INTACT", RGBColor(3, 105, 70))
    _add_chip(slide, 3.05, 2.15, 2.0, "WEAKENING", MIP_AMBER)
    _add_chip(slide, 5.2, 2.15, 2.1, "INVALIDATED", MIP_RED)
    _add_chip(slide, 7.45, 2.15, 2.5, "PROTECTED FULL", MIP_BLUE)

    _add_dark_card(
        slide, 0.7, 2.75, 6.0, 3.7,
        "Symbol Metrics (Live)",
        [
            "Open R:  +1.2R  (reward vs risk from entry)",
            "Expected Move Reached:  84%",
            "Distance to Stop Loss:  1.6%",
            "Distance to Take Profit:  0.8%",
            "Vol Regime:  Normal \u2192 Elevated",
            "Protection Status:  PROTECTED_PARTIAL",
            "",
            "Price chart with TP/SL bands overlaid",
            "Real-time quote updates from IBKR feed",
        ],
        edge=MIP_BLUE,
    )

    _add_dark_card(
        slide, 6.85, 2.75, 5.85, 3.7,
        "Committee Commentary (Live)",
        [
            "Stance:  WATCH_CLOSELY",
            "Confidence:  MEDIUM",
            "",
            "Reason Tags:",
            "  momentum_fade, near_tp, news_uncertainty",
            "",
            "Suggested Actions:",
            "  Tighten stop, monitor revalidation window",
            "  Consider partial exit if thesis weakens further",
        ],
        edge=MIP_PURPLE,
    )

    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(20, 32, 50)
    bar.line.color.rgb = RGBColor(40, 58, 90)
    tf = bar.text_frame
    tf.text = "\U0001F4E1  Like a mission control dashboard: every satellite (position) has live telemetry and a flight director's commentary."
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)

    _add_footer(slide, text="MIP  |  Symbol monitoring + real-time committee commentary")


def _slide_performance(prs: Presentation) -> None:
    """Slide 12: Performance Dashboard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Performance Dashboard",
        "Portfolio-level truth — return, risk, consistency, and the story behind the numbers",
        "REVIEW",
    )

    _add_stat_box(slide, 0.7, 2.15, 2.9, 1.3, "+12.4%", "Total Return", MIP_GREEN)
    _add_stat_box(slide, 3.75, 2.15, 2.9, 1.3, "-4.1%", "Max Drawdown", MIP_RED)
    _add_stat_box(slide, 6.8, 2.15, 2.9, 1.3, "1.82", "Sharpe Ratio", MIP_BLUE)
    _add_stat_box(slide, 9.85, 2.15, 2.8, 1.3, "67%", "Win Rate", MIP_PURPLE)

    _add_card(
        slide, 0.7, 3.7, 4.0, 2.8,
        "Core Metrics",
        [
            "Total equity and return %",
            "Max drawdown and recovery time",
            "Sharpe ratio for risk-adjusted returns",
            "Win rate and expectancy per trade",
            "Monthly cost trend analysis",
        ],
        edge=MIP_GREEN,
    )

    _add_card(
        slide, 4.85, 3.7, 4.0, 2.8,
        "Visual Analysis",
        [
            "Equity curve with drawdown bands",
            "Decision quality trend (expectancy over time)",
            "Selectivity trend — how filters evolve",
            "Decision funnel: signals > proposals > trades",
            "Parallel Worlds counterfactual comparison",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 9.0, 3.7, 3.65, 2.8,
        "Drill-Down Paths",
        [
            "Period filters for focused analysis",
            "Compare return vs drawdown across portfolios",
            "Trace root cause in Training + Decisions",
            "Target realism analysis",
            "Cost attribution breakdown",
        ],
        edge=MIP_PURPLE,
    )

    _add_analogy_bar(
        slide, "\U0001F4CA",
        "Like a financial annual report — but updated daily, with drill-through to every decision that shaped the bottom line.",
    )
    _add_footer(slide)


def _slide_runs(prs: Presentation) -> None:
    """Slide 13: Runs / Audit Trail."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Pipeline Runs & Audit Trail",
        "Operational truth — what ran, when, how long, and what went wrong (if anything)",
        "OPERATIONS",
    )

    _add_chip(slide, 0.7, 2.15, 1.5, "SUCCESS", MIP_GREEN)
    _add_chip(slide, 2.3, 2.15, 2.3, "SUCCESS + SKIPS", MIP_AMBER)
    _add_chip(slide, 4.7, 2.15, 1.3, "FAILED", MIP_RED)
    _add_chip(slide, 6.1, 2.15, 1.5, "RUNNING", MIP_BLUE)

    _add_card(
        slide, 0.7, 2.75, 6.0, 2.1,
        "Daily Pipeline Sequence",
        [
            "Ingest market data \u2192 Generate signals \u2192 Evaluate patterns",
            "Run portfolio simulation \u2192 Propose trades \u2192 Execute approved",
            "Generate briefs + digest \u2192 Run parallel worlds analysis",
            "Scheduled: 5 PM ET, Monday through Friday",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 6.85, 2.75, 5.8, 2.1,
        "Run Detail Panel",
        [
            "Summary cards: status, duration, as-of date, portfolio count",
            "Step-by-step timeline with individual durations",
            "Error panel: SQLSTATE, query IDs, error messages",
            "AI-generated run summary accelerates triage",
        ],
        edge=MIP_PURPLE,
    )

    _add_card(
        slide, 0.7, 5.1, 11.9, 1.4,
        "Audit Log Architecture",
        [
            "MIP_AUDIT_LOG captures every event: stored procedure name, status, duration, error details, and portfolio context.",
            "Immutable audit trail — entries are append-only. Failed runs preserve full diagnostics for root cause analysis.",
            "AI summary explains what happened in plain language — but step-level detail is always available for verification.",
        ],
        edge=MIP_TEAL, item_size=12,
    )

    _add_analogy_bar(
        slide, "\u2708\uFE0F",
        "Like an aircraft's black box recorder: every system event is logged, so investigations start with facts, not guesses.",
    )
    _add_footer(slide)


def _slide_technologies(prs: Presentation) -> None:
    """Slide: Technology Stack."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Technology Stack",
        "Purpose-built on modern, production-grade technologies — Snowflake-native with thin presentation layers",
        "IMPLEMENTATION",
    )

    categories = [
        (
            "Database & AI", MIP_BLUE,
            [
                "Snowflake  —  Data warehouse, compute, and orchestration",
                "Snowflake Cortex  —  Native LLM inference (Claude, Mistral)",
                "Stored Procedures  —  60+ business logic procedures",
                "Scheduled Tasks  —  Daily, hourly, and intraday pipelines",
                "External Access  —  Network rules for market data APIs",
            ],
        ),
        (
            "Backend API", MIP_PURPLE,
            [
                "Python 3.12  —  Runtime for API and scripts",
                "FastAPI  —  High-performance async REST framework",
                "Uvicorn  —  ASGI server for production",
                "snowflake-connector-python  —  Native DB driver",
                "ib-insync  —  Interactive Brokers integration",
            ],
        ),
        (
            "Frontend", MIP_GREEN,
            [
                "React 18  —  Component-based UI framework",
                "Vite 5  —  Lightning-fast build and dev server",
                "Plotly.js + Recharts  —  Interactive financial charts",
                "React Router 6  —  Client-side navigation",
                "react-markdown  —  AI narrative rendering",
            ],
        ),
        (
            "Integrations", MIP_AMBER,
            [
                "Interactive Brokers  —  Live + paper trading",
                "AlphaVantage  —  Historical OHLCV market data",
                "RSS Feeds  —  SEC, Fed, ECB, MarketWatch, IBKR",
                "GitHub Actions  —  CI/CD pipeline",
                "Cursor (Claude)  —  AI-assisted development",
            ],
        ),
    ]

    col_w = 2.95
    gap = 0.1
    for i, (title, color, items) in enumerate(categories):
        x = 0.7 + i * (col_w + gap)
        _add_card(slide, x, 2.1, col_w, 4.4, title, items, edge=color, item_size=10)

    _add_analogy_bar(
        slide, "\U0001F9F1",
        "Like a well-designed building: Snowflake is the foundation, Python is the plumbing, React is the facade, and IBKR is the power grid.",
    )
    _add_footer(slide)


def _slide_platform_architecture(prs: Presentation) -> None:
    """Slide: Platform Architecture — enhanced layered diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, RGBColor(11, 18, 32))
    _draw_header(
        slide,
        "Platform Architecture",
        "Four-layer Snowflake-native architecture — data, logic, and pipelines live together",
        "IMPLEMENTATION",
    )

    dark_bg = RGBColor(15, 23, 42)
    layer_border = RGBColor(36, 48, 65)
    label_color = RGBColor(226, 232, 240)
    desc_color = RGBColor(180, 195, 215)

    layers = [
        ("PRESENTATION LAYER", "React 18 + Vite", [
            "15 feature pages  |  Plotly.js charts  |  Real-time updates",
            "Responsive layout  |  Dark/light theming  |  AI narrative rendering",
        ], MIP_GREEN, 2.05),
        ("API LAYER", "FastAPI + Python 3.12", [
            "25+ REST routers  |  Snowflake connector  |  CORS middleware",
            "IBKR bridge scripts  |  Keypair auth  |  JSON responses",
        ], MIP_PURPLE, 3.25),
        ("LOGIC LAYER", "Snowflake Stored Procedures", [
            "60+ procedures: ingest, signals, training, simulation, proposals, execution",
            "Cortex AI: committee verdicts, digests, narratives  |  Parallel worlds engine",
        ], MIP_BLUE, 4.45),
        ("DATA LAYER", "Snowflake Tables + Views + Tasks", [
            "APP schema: 40+ tables  |  MART schema: analytical views  |  AGENT_OUT: AI output",
            "Scheduled tasks: daily 5PM, intraday 15-min, hourly monitors  |  Audit log",
        ], MIP_TEAL, 5.65),
    ]

    for title, tech, descs, color, y in layers:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(11.9), Inches(1.05)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = dark_bg
        box.line.color.rgb = layer_border

        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(0.1), Inches(1.05)
        )
        edge.fill.solid()
        edge.fill.fore_color.rgb = color
        edge.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = f"{title}    {tech}"
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = label_color
        for d in descs:
            p = tf.add_paragraph()
            p.text = f"    {d}"
            p.font.size = Pt(10)
            p.font.color.rgb = desc_color
            p.space_before = Pt(2)

    arrows_y = [3.12, 4.32, 5.52]
    for ay in arrows_y:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(ay), Inches(0.5), Inches(0.12)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(80, 100, 140)
        arrow.line.fill.background()

    ext_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(2.05), Inches(4.1), Inches(3.0)
    )
    ext_box.fill.solid()
    ext_box.fill.fore_color.rgb = RGBColor(20, 28, 48)
    ext_box.line.color.rgb = MIP_AMBER
    tf = ext_box.text_frame
    tf.word_wrap = True
    tf.text = "EXTERNAL SYSTEMS"
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MIP_AMBER
    externals = [
        "Interactive Brokers (IBKR)",
        "  Paper + Live trading accounts",
        "  Order placement + snapshots",
        "",
        "AlphaVantage",
        "  Historical OHLCV bars",
        "",
        "RSS News Feeds",
        "  SEC, Fed, ECB, MarketWatch",
    ]
    for e in externals:
        p = tf.add_paragraph()
        p.text = f"  {e}"
        p.font.size = Pt(10)
        p.font.color.rgb = desc_color
        p.space_before = Pt(1)

    _add_footer(slide, text="MIP  |  Snowflake-native architecture  |  All logic lives in the warehouse")


def _slide_ai_architecture(prs: Presentation) -> None:
    """Slide: AI Architecture — visual map of every Cortex call in the system."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, RGBColor(11, 18, 32))
    _draw_header(
        slide,
        "AI Architecture \u2014 Snowflake Cortex Integration",
        "8 AI touchpoints across the platform \u2014 every call has a deterministic fallback, every output is auditable",
        "IMPLEMENTATION",
    )

    dark_bg = RGBColor(15, 23, 42)
    card_border = RGBColor(36, 48, 65)
    label_color = RGBColor(226, 232, 240)
    desc_color = RGBColor(180, 195, 215)
    dim_color = RGBColor(130, 145, 170)

    # ── Central Cortex hub ───────────────────────────────────────────────
    hub = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.1), Inches(3.3), Inches(1.2)
    )
    hub.fill.solid()
    hub.fill.fore_color.rgb = MIP_BLUE
    hub.line.fill.background()
    tf = hub.text_frame
    tf.word_wrap = True
    tf.text = "\u2744\uFE0F  SNOWFLAKE CORTEX"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    p = tf.add_paragraph()
    p.text = "LLM inference layer"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.space_before = Pt(2)
    p = tf.add_paragraph()
    p.text = "Claude 3.5 Sonnet  |  Mistral Large 2  |  Llama 3.1 70B"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(170, 195, 240)
    p.space_before = Pt(2)

    # ── AI touchpoint cards — LEFT column (SQL Procedures) ───────────────
    left_calls = [
        ("Daily Digest", "mistral-large2", "Portfolio + global narrative", "SP_AGENT_GENERATE_DAILY_DIGEST", MIP_BLUE),
        ("Training Digest", "mistral-large2", "Maturity journey per symbol", "SP_AGENT_GENERATE_TRAINING_DIGEST", MIP_TEAL),
        ("Committee Verdict", "claude-3.5-sonnet", "Per-proposal trade decision", "SP_VALIDATE_AND_EXECUTE_PROPOSALS", MIP_GREEN),
        ("Parallel Worlds", "mistral-large2", "Counterfactual analysis", "SP_GENERATE_PW_NARRATIVE", MIP_PURPLE),
    ]

    left_x = 0.5
    card_w = 4.3
    card_h = 0.72
    start_y = 2.0
    for i, (name, model, purpose, proc, color) in enumerate(left_calls):
        y = start_y + i * (card_h + 0.12)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_x), Inches(y), Inches(card_w), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = dark_bg
        card.line.color.rgb = card_border

        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left_x), Inches(y), Inches(0.08), Inches(card_h)
        )
        edge.fill.solid()
        edge.fill.fore_color.rgb = color
        edge.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.text = f"  {name}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        p = tf.add_paragraph()
        p.text = f"    {model}  \u2022  {purpose}"
        p.font.size = Pt(9)
        p.font.color.rgb = desc_color
        p.space_before = Pt(1)

        # connector arrow to hub
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(left_x + card_w + 0.06), Inches(y + card_h / 2 - 0.1),
            Inches(0.55), Inches(0.2),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(60, 80, 120)
        arrow.line.fill.background()

    # ── AI touchpoint cards — RIGHT column ───────────────────────────────
    right_calls = [
        ("News Extraction", "llama-3.1-70b", "Structured event parsing", "SP_EXTRACT_NEWS_EVENTS", MIP_RED),
        ("Portfolio Lifecycle", "mistral-large2", "Creation-to-today story", "SP_AGENT_GENERATE_PORTFOLIO_NARRATIVE", MIP_AMBER),
        ("Ask MIP (v2)", "claude-3.5-sonnet", "Route-aware user assistant", "Python API / Orchestrator", MIP_INDIGO),
        ("Ask MIP (v1)", "claude-3.5-sonnet", "Guide-grounded Q&A", "Python API / Legacy Router", RGBColor(113, 128, 150)),
    ]

    right_x = 8.5
    for i, (name, model, purpose, proc, color) in enumerate(right_calls):
        y = start_y + i * (card_h + 0.12)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(right_x), Inches(y), Inches(card_w), Inches(card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = dark_bg
        card.line.color.rgb = card_border

        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(right_x + card_w - 0.08), Inches(y), Inches(0.08), Inches(card_h)
        )
        edge.fill.solid()
        edge.fill.fore_color.rgb = color
        edge.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.text = f"  {name}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        p = tf.add_paragraph()
        p.text = f"    {model}  \u2022  {purpose}"
        p.font.size = Pt(9)
        p.font.color.rgb = desc_color
        p.space_before = Pt(1)

        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(right_x - 0.61), Inches(y + card_h / 2 - 0.1),
            Inches(0.55), Inches(0.2),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(60, 80, 120)
        arrow.line.fill.background()
        arrow.rotation = 180.0

    # ── Section label — left ─────────────────────────────────────────────
    lbl_l = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(4.3), Inches(0.3))
    tf = lbl_l.text_frame
    tf.text = "SQL STORED PROCEDURES"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = dim_color

    lbl_r = slide.shapes.add_textbox(Inches(8.5), Inches(5.5), Inches(4.3), Inches(0.3))
    tf = lbl_r.text_frame
    tf.text = "PYTHON API  +  SQL PROCEDURES"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = dim_color

    # ── Fallback banner ──────────────────────────────────────────────────
    fb = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.7)
    )
    fb.fill.solid()
    fb.fill.fore_color.rgb = RGBColor(20, 32, 55)
    fb.line.color.rgb = RGBColor(40, 58, 90)
    tf = fb.text_frame
    tf.word_wrap = True
    tf.text = "SAFETY NET  \u2014  Every AI call has a deterministic fallback"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MIP_AMBER
    p = tf.add_paragraph()
    p.text = (
        "    If Cortex fails or times out: digests generate bullet-point summaries from snapshots, "
        "committee falls back to deterministic approval, news uses heuristic extraction, "
        "and Ask MIP returns a service error.  No AI failure can halt the trading pipeline."
    )
    p.font.size = Pt(10)
    p.font.color.rgb = desc_color
    p.space_before = Pt(2)

    # ── Output targets ───────────────────────────────────────────────────
    out_label = slide.shapes.add_textbox(Inches(5.0), Inches(4.5), Inches(3.3), Inches(0.3))
    tf = out_label.text_frame
    tf.text = "Output \u2192 AGENT_OUT schema"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MIP_GREEN

    out_down = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(4.35), Inches(0.5), Inches(0.18)
    )
    out_down.fill.solid()
    out_down.fill.fore_color.rgb = MIP_GREEN
    out_down.line.fill.background()

    tables = slide.shapes.add_textbox(Inches(4.2), Inches(4.8), Inches(4.8), Inches(0.6))
    tf = tables.text_frame
    tf.word_wrap = True
    tf.text = (
        "DAILY_DIGEST_NARRATIVE  \u2022  TRAINING_DIGEST_NARRATIVE  \u2022  "
        "PORTFOLIO_LIFECYCLE_NARRATIVE  \u2022  PARALLEL_WORLD_NARRATIVE  \u2022  "
        "ORDER_PROPOSALS  \u2022  NEWS_EVENT_EXTRACTED  \u2022  ASK_QUERY_EVENT"
    )
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(8)
    tf.paragraphs[0].font.color.rgb = dim_color

    _add_footer(slide, text="MIP  |  8 AI touchpoints  |  3 models  |  Full deterministic fallback")


def _slide_ask_mip(prs: Presentation) -> None:
    """Slide: Ask MIP — the intelligent route-aware assistant."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_BG)
    _draw_header(
        slide,
        "Ask MIP \u2014 Your Intelligent Platform Assistant",
        "Route-aware, docs-grounded AI assistant — always one click away on every page",
        "FEATURE",
    )

    # ── Flow diagram ─────────────────────────────────────────────────────
    _add_process(
        slide,
        ["User Question", "Route Context", "Doc Retrieval", "Glossary Match", "Cortex LLM", "Answer"],
        y=2.1, start_x=0.4, box_w=1.78, gap=0.3,
    )

    # ── Three feature columns ────────────────────────────────────────────
    _add_card(
        slide, 0.7, 3.25, 3.9, 3.2,
        "How It Works",
        [
            "Floating button on every page (\u2753 FAB)",
            "Panel slides open from the right",
            "User types a question in natural language",
            "System detects current route + page context",
            "Retrieves relevant docs + glossary matches",
            "Calls Snowflake Cortex (Claude 3.5 Sonnet)",
            "Returns structured answer with provenance",
        ],
        edge=MIP_BLUE,
    )

    _add_card(
        slide, 4.75, 3.25, 4.0, 3.2,
        "Route-Aware Intelligence",
        [
            "Knows which page you're viewing",
            "Injects page-specific guide content as context",
            "Frames answers relative to visible UI elements",
            "24 route-to-guide mappings for precise context",
            "Score boosting for current-page documentation",
            "Page hint block in system prompt (first 1200 chars)",
            "General finance knowledge when MIP docs don't cover",
        ],
        edge=MIP_PURPLE,
    )

    _add_card(
        slide, 8.9, 3.25, 3.75, 3.2,
        "Quality & Trust",
        [
            "Provenance badges: DOC, GLOSSARY, INFERENCE",
            "Confidence scoring (docs 50%, glossary 35%)",
            "\"Did you mean\" suggestions when unsure",
            "Unknown terms logged for glossary review",
            "Telemetry: every query tracked for coverage",
            "Never invents live values or thresholds",
            "90-second timeout, last 10 messages as context",
        ],
        edge=MIP_GREEN,
    )

    # ── Architecture strip ───────────────────────────────────────────────
    _add_chip(slide, 0.7, 6.55, 2.3, "CLAUDE 3.5 SONNET", MIP_INDIGO)
    _add_chip(slide, 3.1, 6.55, 2.5, "24 GUIDE FILES", MIP_BLUE)
    _add_chip(slide, 5.7, 6.55, 2.3, "GLOSSARY TABLE", MIP_TEAL)
    _add_chip(slide, 8.1, 6.55, 2.2, "TELEMETRY LOG", MIP_PURPLE)
    _add_chip(slide, 10.4, 6.55, 2.2, "FALLBACK SAFE", MIP_GREEN)

    _add_analogy_bar(
        slide, "\U0001F9ED",
        "Like a knowledgeable colleague sitting next to you: they know the platform, see your screen, and answer questions grounded in facts \u2014 not guesses.",
    )
    _add_footer(slide)


def _slide_how_mip_was_built(prs: Presentation) -> None:
    """Slide: How MIP Was Built — The AI-Orchestrated Story."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_DARK)

    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = MIP_DARK
    banner.line.fill.background()

    # ── Title ────────────────────────────────────────────────────────────
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(12.0), Inches(0.7))
    tf = title.text_frame
    tf.text = "How MIP Was Built \u2014 The AI-Orchestrated Story"
    tf.paragraphs[0].font.size = Pt(30)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.04)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = MIP_BLUE
    accent_line.line.fill.background()

    # ── Subtitle ─────────────────────────────────────────────────────────
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(10.0), Inches(0.35))
    tf = sub.text_frame
    tf.text = "Development Timeline \u2014 6 months, zero lines of code typed by a human"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)

    # ── Timeline bar ─────────────────────────────────────────────────────
    timeline_y = 1.65
    timeline_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(timeline_y + 0.22), Inches(11.9), Inches(0.06)
    )
    timeline_bar.fill.solid()
    timeline_bar.fill.fore_color.rgb = RGBColor(40, 58, 90)
    timeline_bar.line.fill.background()

    milestones = [
        (0.7, "Oct \u2013 Dec 2025"),
        (4.7, "Jan 2026"),
        (8.2, "Feb 2026"),
        (10.8, "Mar 2026"),
    ]
    for mx, mlabel in milestones:
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(mx + 0.3), Inches(timeline_y + 0.12), Inches(0.26), Inches(0.26)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = MIP_BLUE
        dot.line.fill.background()
        t = slide.shapes.add_textbox(Inches(mx), Inches(timeline_y - 0.15), Inches(1.8), Inches(0.25))
        tf = t.text_frame
        tf.text = mlabel
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = MIP_BLUE

    # ── Phase cards ──────────────────────────────────────────────────────
    phase_y = 2.15
    phase_h = 1.55
    phases = [
        (0.7, 2.85, "ChatGPT + Codex", MIP_BLUE, [
            "ChatGPT as investment SME",
            "& solution architect.",
            "OpenAI Codex as developer.",
            "Core pipeline built.",
        ]),
        (3.65, 2.85, "Transition to Cursor", MIP_PURPLE, [
            "Switched to Cursor (Claude)",
            "for development & implementation.",
            "ChatGPT for requirements",
            "& idea validation.",
        ]),
        (6.6, 2.85, "Full Cursor Autonomy", MIP_GREEN, [
            "Cursor deploys direct to",
            "Snowflake. Testing, debugging,",
            "data repairs \u2014 all autonomous.",
            "Human reviews, never types.",
        ]),
        (9.55, 2.85, "Platform Maturity", MIP_AMBER, [
            "IBKR live broker integration.",
            "Cortex AI committee + narratives.",
            "Parallel worlds + news intelligence.",
            "146K lines. Production-ready.",
        ]),
    ]
    for px, pw, ptitle, pcolor, plines in phases:
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(phase_y), Inches(pw), Inches(phase_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(20, 32, 55)
        card.line.color.rgb = pcolor

        edge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(px), Inches(phase_y), Inches(pw), Inches(0.06)
        )
        edge.fill.solid()
        edge.fill.fore_color.rgb = pcolor
        edge.line.fill.background()

        tf = card.text_frame
        tf.word_wrap = True
        tf.text = ptitle
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = pcolor
        for ln in plines:
            p = tf.add_paragraph()
            p.text = ln
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(180, 195, 215)
            p.space_before = Pt(2)

    # ── Kenneth's Role section ───────────────────────────────────────────
    role_y = 4.0
    role_title = slide.shapes.add_textbox(Inches(0.7), Inches(role_y), Inches(5.0), Inches(0.35))
    tf = role_title.text_frame
    tf.text = "Kenneth's Role"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)

    roles = [
        ("\u2B50", "Visionary", "Conceived the mission.\nDefined what MIP should\nbecome.", MIP_AMBER),
        ("\u2699\uFE0F", "Orchestrator", "Directed the AI agents.\nReviewed output, steered\ndirection.", MIP_BLUE),
        ("\u2611\uFE0F", "Product Owner", "Prioritised features.\nDefined acceptance criteria.\nQuality gatekeeper.", MIP_GREEN),
    ]
    for i, (icon, rname, rdesc, rcolor) in enumerate(roles):
        rx = 0.7 + i * 2.15
        rcard = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(rx), Inches(role_y + 0.4), Inches(2.0), Inches(1.4)
        )
        rcard.fill.solid()
        rcard.fill.fore_color.rgb = RGBColor(20, 32, 55)
        rcard.line.color.rgb = RGBColor(40, 58, 90)
        tf = rcard.text_frame
        tf.word_wrap = True
        tf.text = f"{icon}  {rname}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = rcolor
        for dline in rdesc.split("\n"):
            p = tf.add_paragraph()
            p.text = dline
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(160, 175, 195)
            p.space_before = Pt(1)

    # ── The Numbers table ────────────────────────────────────────────────
    tbl_title = slide.shapes.add_textbox(Inches(7.0), Inches(role_y), Inches(5.0), Inches(0.35))
    tf = tbl_title.text_frame
    tf.text = "The Numbers"
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(226, 232, 240)

    table_data = [
        ("SQL (Snowflake)", "340 files", "47,500 lines", "40%", MIP_AMBER),
        ("React (JSX/JS/CSS)", "100 files", "31,500 lines", "27%", MIP_BLUE),
        ("Python (FastAPI)", "51 files", "23,100 lines", "20%", MIP_GREEN),
        ("Python (Scripts)", "23 files", "10,300 lines", "9%", MIP_PURPLE),
        ("Documentation", "80 files", "5,300 lines", "4%", RGBColor(113, 128, 150)),
    ]

    tbl_x = 7.0
    tbl_y = role_y + 0.42
    row_h = 0.26
    col_widths = [2.5, 1.0, 1.3, 0.7]

    header_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(tbl_x), Inches(tbl_y), Inches(5.6), Inches(row_h)
    )
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = RGBColor(30, 42, 65)
    header_bg.line.fill.background()
    headers = ["Layer", "Files", "Lines", "%"]
    for ci, (htext, cw) in enumerate(zip(headers, col_widths)):
        hx = tbl_x + sum(col_widths[:ci]) + 0.1
        ht = slide.shapes.add_textbox(Inches(hx), Inches(tbl_y), Inches(cw), Inches(row_h))
        tf = ht.text_frame
        tf.text = htext
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(180, 195, 215)

    for ri, (layer, files, lines, pct, dot_color) in enumerate(table_data):
        ry = tbl_y + (ri + 1) * row_h
        if ri % 2 == 0:
            row_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(tbl_x), Inches(ry), Inches(5.6), Inches(row_h)
            )
            row_bg.fill.solid()
            row_bg.fill.fore_color.rgb = RGBColor(20, 30, 50)
            row_bg.line.fill.background()

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(tbl_x + 0.05), Inches(ry + 0.07), Inches(0.12), Inches(0.12)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = dot_color
        dot.line.fill.background()

        vals = [layer, files, lines, pct]
        for ci, (val, cw) in enumerate(zip(vals, col_widths)):
            vx = tbl_x + sum(col_widths[:ci]) + 0.1
            vt = slide.shapes.add_textbox(Inches(vx + (0.12 if ci == 0 else 0)), Inches(ry), Inches(cw), Inches(row_h))
            tf = vt.text_frame
            tf.text = val
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = RGBColor(210, 220, 235)

    total_y = tbl_y + 6 * row_h
    total_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(tbl_x), Inches(total_y), Inches(5.6), Inches(row_h + 0.02)
    )
    total_bg.fill.solid()
    total_bg.fill.fore_color.rgb = MIP_BLUE
    total_bg.line.fill.background()
    totals = ["TOTAL", "594", "117,700", "100%"]
    for ci, (val, cw) in enumerate(zip(totals, col_widths)):
        vx = tbl_x + sum(col_widths[:ci]) + 0.1
        vt = slide.shapes.add_textbox(Inches(vx), Inches(total_y), Inches(cw), Inches(row_h + 0.02))
        tf = vt.text_frame
        tf.text = val
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE

    # ── Quote banner ─────────────────────────────────────────────────────
    quote_y = 6.3
    quote_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(quote_y), Inches(11.9), Inches(0.95)
    )
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = RGBColor(20, 35, 65)
    quote_bg.line.color.rgb = MIP_BLUE

    qt = slide.shapes.add_textbox(Inches(1.0), Inches(quote_y + 0.08), Inches(11.4), Inches(0.8))
    tf = qt.text_frame
    tf.word_wrap = True
    tf.text = (
        "\u201CKenneth never wrote a line of code, AI did.  "
        "AI never made a product decision, Kenneth did.\n"
        "Each played to their strengths.  117,700 lines.  594 files.  "
        "6 months.  One human involved.\u201D"
    )
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = RGBColor(220, 230, 245)
    if len(tf.paragraphs) > 1:
        tf.paragraphs[1].font.size = Pt(14)
        tf.paragraphs[1].font.italic = True
        tf.paragraphs[1].font.color.rgb = RGBColor(220, 230, 245)


def _slide_closing(prs: Presentation) -> None:
    """Slide 15: Closing / CTA."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, MIP_DARK)

    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = MIP_DARK
    banner.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(1.0), Inches(0.15), Inches(5.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = MIP_BLUE
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(1.3), Inches(1.2), Inches(10.5), Inches(1.2))
    tf = title.text_frame
    tf.text = "MIP: Intelligence Meets Execution"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(241, 245, 249)

    sub = slide.shapes.add_textbox(Inches(1.3), Inches(2.5), Inches(10.0), Inches(0.6))
    tf = sub.text_frame
    tf.text = "The platform where every trade is born from evidence, governed by committee, and remembered for the future."
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(3.3), Inches(4.0), Inches(0.04)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = MIP_BLUE
    divider.line.fill.background()

    pillars_data = [
        ("\U0001F52C", "Research", "Pattern detection, training maturity, news intelligence, parallel worlds"),
        ("\u2696\uFE0F", "Decisions", "AI committee verdicts, reason codes, structured gates, full transparency"),
        ("\U0001F4B9", "Execution", "Live IBKR link, position lifecycle, risk brakes, real-time sync"),
        ("\U0001F4DA", "Learning", "Causal ledger, outcome tracking, policy feedback, continuous improvement"),
    ]

    y = 3.6
    for icon, pillar, desc in pillars_data:
        row = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(10.0), Inches(0.55))
        tf = row.text_frame
        tf.word_wrap = True
        tf.text = f"{icon}  {pillar}  —  {desc}"
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = RGBColor(203, 213, 225)
        y += 0.65

    cta = slide.shapes.add_textbox(Inches(1.3), Inches(6.2), Inches(10.0), Inches(0.5))
    tf = cta.text_frame
    tf.text = "Ready to see it live?  Let's walk through the platform together."
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = MIP_BLUE


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD DECK
# ══════════════════════════════════════════════════════════════════════════════

def build_deck(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_title(prs)               # 1  - Title / Hero
    _slide_what_is_mip(prs)         # 2  - What is MIP
    _slide_training_status(prs)     # 3  - Training Status
    _slide_market_timeline(prs)     # 4  - Market Timeline
    _slide_news_intelligence(prs)   # 5  - News Intelligence
    _slide_parallel_worlds(prs)     # 6  - Parallel Worlds
    _slide_live_link(prs)           # 7  - Live Portfolio Link
    _slide_live_activity(prs)       # 8  - Live Portfolio Activity
    _slide_ai_decisions(prs)        # 9  - AI Agent Decisions
    _slide_learning_ledger(prs)     # 10 - Learning Ledger
    _slide_symbol_tracker(prs)      # 11 - Live Symbol Tracker
    _slide_performance(prs)         # 12 - Performance
    _slide_runs(prs)                # 13 - Runs / Audit
    _slide_ask_mip(prs)             # 14 - Ask MIP
    _slide_technologies(prs)        # 15 - Technology Stack
    _slide_platform_architecture(prs)  # 16 - Architecture Diagram
    _slide_ai_architecture(prs)     # 17 - AI Architecture (Cortex calls)
    _slide_how_mip_was_built(prs)   # 18 - How MIP Was Built
    _slide_closing(prs)             # 19 - Closing / CTA

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    deck_path = Path(__file__).with_name("MIP_Platform_Introduction.pptx")
    build_deck(deck_path)
    print(f"Created: {deck_path}")
