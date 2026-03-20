from pathlib import Path

from pptx import Presentation


def set_text(shape, text: str) -> None:
    tf = shape.text_frame
    tf.clear()
    lines = text.split("\n")
    tf.paragraphs[0].text = lines[0] if lines else ""
    for line in lines[1:]:
        tf.add_paragraph().text = line


def clear_slide_text(slide) -> None:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()


def apply_content(prs: Presentation) -> None:
    # Slide 1 - Title
    s = prs.slides[0]
    clear_slide_text(s)
    set_text(s.shapes[23], "Market Intelligence Platform (MIP)")
    set_text(
        s.shapes[25],
        "Research Internals to Trading Operations\nDemo + Training for Colleagues",
    )
    set_text(s.shapes[26], "Kenneth Jebjerg    March 2026    Internal Training")

    # Slide 2 - Agenda sequence
    s = prs.slides[1]
    clear_slide_text(s)
    set_text(s.shapes[1], "Inside-Out Demo Sequence")
    set_text(s.shapes[6], "Research\nFoundation")
    set_text(s.shapes[7], "MIP intro, training status,\nmarket timeline, parallel worlds,\nnews intelligence.")
    set_text(s.shapes[11], "Decision\nLayer")
    set_text(s.shapes[12], "AI Agent Decisions and\ncommittee reasoning become\naction guidance.")
    set_text(s.shapes[16], "Live Trading\nOperations")
    set_text(s.shapes[17], "Live portfolio link + activity,\nlive symbol tracker, committee\ncommentary in context.")
    set_text(s.shapes[21], "Review,\nLearning, Audit")
    set_text(s.shapes[22], "Learning ledger, performance,\naudit trail and operator\nplaybook.")
    set_text(s.shapes[24], "Every section links evidence -> committee -> action -> review.")

    # Slide 3 - Intro pipeline
    s = prs.slides[2]
    clear_slide_text(s)
    set_text(s.shapes[1], "MIP Introduction")
    set_text(
        s.shapes[3],
        "MIP is an evidence-first decision system: learn from repeatable behavior,\nthen propose actions through deterministic controls.",
    )
    set_text(s.shapes[4], "Ingest")
    set_text(s.shapes[6], "Detect")
    set_text(s.shapes[8], "Evaluate")
    set_text(s.shapes[10], "Trust")
    set_text(s.shapes[12], "Decide")
    set_text(s.shapes[14], "Operate")
    set_text(s.shapes[15], "Market bars\n(OHLC)")
    set_text(s.shapes[16], "Pattern\nsignals")
    set_text(s.shapes[17], "H1/H3/H5/H10/H20\noutcomes")
    set_text(s.shapes[18], "Maturity + coverage +\ntrust labels")
    set_text(s.shapes[19], "Committee verdicts\nand reason codes")
    set_text(s.shapes[20], "Live workflow,\nperformance, audit")
    set_text(
        s.shapes[22],
        "Research and trading are one continuous loop, not separate systems.",
    )

    # Slide 4 - Navigation map
    s = prs.slides[3]
    clear_slide_text(s)
    set_text(s.shapes[1], "Current MIP Operating Map")
    set_text(
        s.shapes[3],
        "Updated navigation and responsibilities across pages in the current UI.",
    )
    set_text(s.shapes[5], "Dashboard")
    set_text(s.shapes[6], "Cockpit, Home,\nPerformance")
    set_text(s.shapes[8], "Portfolio")
    set_text(s.shapes[10], "Live Portfolio Link,\nLive Activity,\nLive Symbol Tracker")
    set_text(s.shapes[12], "Research")
    set_text(s.shapes[14], "Training Status,\nMarket Timeline,\nParallel Worlds,\nNews Intelligence")
    set_text(s.shapes[16], "Decision Executions")
    set_text(s.shapes[17], "AI Agent Decisions,\nLearning Ledger")
    set_text(s.shapes[19], "Operations")
    set_text(s.shapes[21], "Runs (Audit),\nDebug")
    set_text(s.shapes[23], "Ask MIP")
    set_text(s.shapes[27], "Cross-page AI assistant\nfor plain-language guidance")
    set_text(s.shapes[28], "Evidence + decisions + operations are connected in one flow.")
    set_text(s.shapes[29], "APP")
    set_text(s.shapes[30], "MART")
    set_text(s.shapes[31], "AGENT_OUT")
    set_text(s.shapes[32], "RAW_EXT")
    set_text(s.shapes[34], "Core Data Zones")
    set_text(s.shapes[35], "Execution tables,\ntrades, controls")
    set_text(s.shapes[37], "Analytics views,\nKPIs, training")
    set_text(s.shapes[40], "Narratives,\ndigests,\ncommittee output")
    set_text(s.shapes[43], "External market +\nnews ingestion")
    set_text(s.shapes[46], "Why this matters:")
    set_text(
        s.shapes[47],
        "One platform model: deterministic data, bounded AI narration, and full auditability.",
    )

    # Slide 5 - Training status
    s = prs.slides[4]
    clear_slide_text(s)
    set_text(s.shapes[1], "Training Status")
    set_text(
        s.shapes[3],
        "Core question: how much evidence do we have for each symbol + pattern before trading influence?",
    )
    set_text(s.shapes[4], "How to Read")
    set_text(
        s.shapes[6],
        "Start with sample size + coverage.\nThen evaluate horizon averages.\nMaturity is confidence-in-evidence.",
    )
    set_text(s.shapes[7], "Maturity Bands")
    set_text(s.shapes[10], "TRUST READY")
    set_text(s.shapes[11], "High sample + coverage,\nconsistent outcomes")
    set_text(s.shapes[12], "Eligible for proposal flow")
    set_text(s.shapes[14], "LEARNING")
    set_text(s.shapes[15], "Growing evidence but\nnot fully stable yet")
    set_text(s.shapes[16], "Monitor before promoting")
    set_text(s.shapes[18], "EARLY")
    set_text(s.shapes[19], "Insufficient observations\nor weak coverage")
    set_text(s.shapes[20], "Research only")
    set_text(s.shapes[22], "Tracked horizons:")
    set_text(s.shapes[24], "H1")
    set_text(s.shapes[26], "H3")
    set_text(s.shapes[28], "H5")
    set_text(s.shapes[30], "H10")
    set_text(s.shapes[32], "H20")
    set_text(s.shapes[33], "Example Interpretation")
    set_text(s.shapes[35], "Avg H5 = +0.8%")
    set_text(s.shapes[38], "Coverage = 86%")
    set_text(s.shapes[41], "Maturity = 78")
    set_text(s.shapes[44], "AI note: Ask MIP can translate rows into plain-language confidence.")

    # Slide 6 - Market timeline
    s = prs.slides[5]
    clear_slide_text(s)
    set_text(s.shapes[1], "Market Timeline")
    set_text(
        s.shapes[3],
        "Symbol-level observability from signal to proposal to execution, with narrative context.",
    )
    set_text(s.shapes[4], "Timeline Funnel")
    set_text(s.shapes[6], "Signals")
    set_text(s.shapes[7], "Detections in selected window")
    set_text(s.shapes[10], "Proposals")
    set_text(s.shapes[11], "Signals passing trust/risk filters")
    set_text(s.shapes[14], "Trades")
    set_text(s.shapes[15], "Executed actions by portfolio")
    set_text(s.shapes[18], "Narrative + Chain")
    set_text(s.shapes[19], "Decision narrative explains why flow advanced or stalled")
    set_text(s.shapes[21], "Tile states: executed, proposed, signal-only, inactive")
    set_text(s.shapes[23], "Cross-links: Live Activity, AI Agent Decisions, Runs")
    set_text(s.shapes[26], "Concrete Example")
    set_text(s.shapes[27], "Symbol")
    set_text(s.shapes[28], "AAPL")
    set_text(s.shapes[29], "Signals")
    set_text(s.shapes[30], "12")
    set_text(s.shapes[31], "Proposals")
    set_text(s.shapes[32], "3")
    set_text(s.shapes[33], "Trades")
    set_text(s.shapes[34], "1")
    set_text(s.shapes[35], "Trust")
    set_text(s.shapes[36], "TRUSTED")
    set_text(s.shapes[37], "Open Status")
    set_text(s.shapes[38], "WATCH CLOSELY")
    set_text(
        s.shapes[47],
        "Use this page to tell the symbol story end-to-end in the demo.",
    )

    # Slide 7 - Live portfolio link
    s = prs.slides[6]
    clear_slide_text(s)
    set_text(s.shapes[1], "Live Portfolio Link")
    set_text(s.shapes[3], "Control plane connecting MIP live workflow to broker truth.")
    set_text(s.shapes[6], "Create/edit live config with system-assigned IDs")
    set_text(s.shapes[8], "Bind IBKR account and adapter mode")
    set_text(s.shapes[10], "Set freshness, size, and drawdown controls")
    set_text(s.shapes[12], "Validate activation guard and readiness chain")
    set_text(s.shapes[14], "Research source selected at import in Live Activity")
    set_text(s.shapes[16], "Saving config writes governance state only")
    set_text(s.shapes[17], "Execution Readiness States")
    set_text(s.shapes[20], "CONNECTED")
    set_text(s.shapes[21], "Wiring and config pass checks")
    set_text(s.shapes[23], "CAUTION")
    set_text(s.shapes[24], "Freshness/drift or limit threshold approaching")
    set_text(s.shapes[26], "BLOCKED")
    set_text(s.shapes[27], "Readiness failed; new execution blocked")
    set_text(s.shapes[28], "Key control fields:")
    set_text(s.shapes[29], "Max Positions")
    set_text(s.shapes[30], "Max concurrent holdings")
    set_text(s.shapes[31], "Max Position %")
    set_text(s.shapes[32], "Per-position exposure cap")
    set_text(s.shapes[33], "Quote/Snapshot Freshness")
    set_text(s.shapes[34], "Stale-data blocker")
    set_text(s.shapes[35], "Drawdown Stop / Bust %")
    set_text(s.shapes[36], "Portfolio safety brakes")
    set_text(s.shapes[43], "Source -> Live Portfolio -> IBKR -> Guard -> Readiness")
    set_text(s.shapes[44], "Full control chain:")

    # Slide 8 - Live portfolio activity
    s = prs.slides[7]
    clear_slide_text(s)
    set_text(s.shapes[1], "Live Portfolio Activity")
    set_text(
        s.shapes[3],
        "Operational lifecycle for live-linked paper workflow with committee checkpoints.",
    )
    set_text(s.shapes[5], "IMPORTED")
    set_text(s.shapes[6], "Proposal imported\ninto live flow")
    set_text(s.shapes[8], "VALIDATED")
    set_text(s.shapes[9], "Freshness + realism\nchecks passed")
    set_text(s.shapes[11], "COMMITTEE")
    set_text(s.shapes[12], "Verdict + reason codes\nrecorded")
    set_text(s.shapes[14], "APPROVED")
    set_text(s.shapes[15], "PM/compliance +\nrevalidation success")
    set_text(s.shapes[20], "EXECUTED")
    set_text(s.shapes[21], "Action completed\nand tracked")
    set_text(s.shapes[22], "Live Activity Use")
    set_text(s.shapes[23], "Confirm freshness and latest transitions")
    set_text(s.shapes[25], "Investigate delays via reason fields")
    set_text(
        s.shapes[27],
        "Cross-check with AI Agent Decisions + Runs.\nCommittee touchpoint is explicit and auditable.",
    )

    # Slide 9 - News Intelligence
    s = prs.slides[8]
    clear_slide_text(s)
    set_text(s.shapes[1], "News Intelligence")
    set_text(
        s.shapes[3],
        "Evidence-backed news context for decisions. Explainable and bounded, not narrative guesswork.",
    )
    set_text(s.shapes[7], "CONTEXT")
    set_text(s.shapes[8], "Market KPIs")
    set_text(s.shapes[9], "Symbols with news\nHOT symbols\nsnapshot age")
    set_text(s.shapes[14], "IMPACT")
    set_text(s.shapes[15], "Decision Impact")
    set_text(s.shapes[16], "Proposals scoped\nwith context/adj\nblocked by news risk")
    set_text(s.shapes[21], "HOT")
    set_text(s.shapes[22], "Sidebar marker")
    set_text(s.shapes[23], "Unseen + decision-relevant\nnews only")
    set_text(s.shapes[28], "GUARDRAILS")
    set_text(s.shapes[29], "Quality controls")
    set_text(s.shapes[30], "Invalid URLs excluded\nfreshness explicit\ninfluence bounded")

    # Slide 10 - AI Agent Decisions
    s = prs.slides[9]
    clear_slide_text(s)
    set_text(s.shapes[1], "AI Agent Decisions (Committee)")
    set_text(
        s.shapes[3],
        "Decision courtroom for simulation + live workflows: verdict, reason codes, and revalidation.",
    )
    set_text(s.shapes[6], "1. Proposal\nContext")
    set_text(
        s.shapes[7],
        "Structured snapshot:\ntrust, risk, capacity,\nnews, freshness, timing.",
    )
    set_text(s.shapes[11], "2. Committee\nEvaluation")
    set_text(
        s.shapes[12],
        "Verdict + summary +\nreason tags generated\nfrom bounded evidence.",
    )
    set_text(s.shapes[16], "3. Revalidation\nBefore Action")
    set_text(
        s.shapes[17],
        "Final checks before execution:\nif conditions changed,\nstatus can downgrade.",
    )
    set_text(s.shapes[19], "Typical status journey: PROPOSED -> APPROVED -> REJECTED or EXECUTED")

    # Slide 11 - Parallel Worlds
    s = prs.slides[10]
    clear_slide_text(s)
    set_text(s.shapes[1], "Parallel Worlds")
    set_text(
        s.shapes[3],
        "Counterfactual lab: replay same market data through alternative rules and measure regret safely.",
    )
    set_text(s.shapes[4], "What Is Tested")
    set_text(s.shapes[7], "Signal Filter")
    set_text(s.shapes[8], "Looser/tighter threshold\nfor proposal eligibility")
    set_text(s.shapes[10], "Example Insight")
    set_text(s.shapes[11], "Looser filter improves PnL today,\nbut confidence = Noise.")
    set_text(s.shapes[14], "Position Size")
    set_text(s.shapes[15], "50%-150% replay\nof current sizing")
    set_text(s.shapes[17], "Example Insight")
    set_text(s.shapes[18], "125% size increased return\nwith higher drawdown.")
    set_text(s.shapes[21], "Entry Timing")
    set_text(s.shapes[22], "Delay 0-3 bars\nfor execution timing")
    set_text(s.shapes[24], "Example Insight")
    set_text(s.shapes[25], "Delay 1 bar reduced regret\nacross recent regime.")
    set_text(s.shapes[28], "Baseline")
    set_text(s.shapes[29], "Stay in cash\nor no-rule benchmark")
    set_text(s.shapes[31], "Example Insight")
    set_text(s.shapes[32], "Cash baseline highlights\nwhen risk controls saved equity.")
    set_text(s.shapes[35], "Policy Health outputs:")
    set_text(s.shapes[41], "- Regret attribution by category")
    set_text(s.shapes[42], "- Scenario comparison and confidence tiers")
    set_text(s.shapes[43], "- Equity overlays + AI analysis")
    set_text(s.shapes[44], "- Safety framing: analysis only, no auto-apply")

    # Slide 12 - Live symbol tracker + committee commentary
    s = prs.slides[11]
    clear_slide_text(s)
    set_text(s.shapes[0], "Live Symbol Tracker")
    set_text(s.shapes[7], "Symbol State")
    set_text(s.shapes[12], "Thesis labels:\nTHESIS_INTACT,\nWEAKENING,\nINVALIDATED")
    set_text(s.shapes[18], "Risk Math")
    set_text(s.shapes[23], "Open R, expected move reached,\ndistance to TP/SL,\nposition protection")
    set_text(s.shapes[29], "Committee Commentary")
    set_text(s.shapes[34], "Stance + confidence +\nreason tags + actions to consider")
    set_text(s.shapes[40], "Cross-check Path")
    set_text(s.shapes[45], "Symbol Tracker -> AI Agent Decisions -> News Intelligence -> Runs")
    set_text(s.shapes[51], "Live Monitoring")
    set_text(s.shapes[56], "Use during intraday windows for elevated symbols first.")
    set_text(s.shapes[62], "Operator Rule")
    set_text(s.shapes[67], "Commentary is guidance; deterministic gates still control execution.")
    set_text(s.shapes[73], "Outcome Link")
    set_text(s.shapes[78], "Trace outcomes later in Learning Ledger for feedback.")

    # Slide 13 - AI naturally fused + governance
    s = prs.slides[12]
    clear_slide_text(s)
    set_text(s.shapes[1], "How AI Is Naturally Fused")
    set_text(s.shapes[4], "AI Helps")
    set_text(s.shapes[8], "Committee summaries and verdict context")
    set_text(s.shapes[11], "Parallel Worlds narrative and regret interpretation")
    set_text(s.shapes[14], "News reader summary and HOT prioritization")
    set_text(s.shapes[17], "Ask MIP route-aware coaching")
    set_text(s.shapes[25], "Humans Stay In Control")
    set_text(s.shapes[29], "Risk gates and thresholds are deterministic")
    set_text(s.shapes[32], "Execution requires approvals and revalidation")
    set_text(s.shapes[35], "AI outputs are non-binding explanations")
    set_text(s.shapes[38], "Runs + Ledger provide accountability trail")

    # Slide 14 - Learning Ledger
    s = prs.slides[13]
    clear_slide_text(s)
    set_text(s.shapes[1], "Learning Ledger")
    set_text(s.shapes[6], "Decision chain from proposal to realized outcome")
    set_text(s.shapes[11], "Evidence attribution by role and context")
    set_text(s.shapes[16], "Expected vs realized outcome comparison")
    set_text(s.shapes[21], "Behavior change review across periods")
    set_text(s.shapes[26], "Post-trade review")
    set_text(s.shapes[27], "Explain why outcomes improved or degraded")
    set_text(s.shapes[31], "Weekly retrospective")
    set_text(s.shapes[32], "Find repeat rejection or drift patterns")
    set_text(s.shapes[31], "Governance")
    set_text(s.shapes[32], "Support stakeholder and compliance explainability")

    # Slide 15 - Performance
    s = prs.slides[14]
    clear_slide_text(s)
    set_text(s.shapes[22], "Performance Dashboard")
    set_text(
        s.shapes[24],
        "Portfolio-level truth: compare return, drawdown, and consistency before changing policy.",
    )
    set_text(s.shapes[26], "Level 1")
    set_text(s.shapes[27], "Snapshot")
    set_text(s.shapes[28], "Current return and drawdown\nby portfolio")
    set_text(s.shapes[29], "Live")
    set_text(s.shapes[31], "Level 2")
    set_text(s.shapes[32], "Trend")
    set_text(s.shapes[33], "Period filters (1M/3M/YTD)\nfor direction and stability")
    set_text(s.shapes[34], "Live")
    set_text(s.shapes[36], "Level 3")
    set_text(s.shapes[37], "Attribution")
    set_text(s.shapes[38], "Link movement to training,\ndecisions, and risk gates")
    set_text(s.shapes[39], "Live")
    set_text(s.shapes[41], "Level 4")
    set_text(s.shapes[42], "Policy")
    set_text(s.shapes[43], "Use Parallel Worlds and\ncommittee logs for tuning")
    set_text(s.shapes[44], "In Use")
    set_text(s.shapes[46], "Level 5")
    set_text(s.shapes[47], "Continuous Learning")
    set_text(s.shapes[48], "Ledger feedback loops improve\nfuture committee behavior")
    set_text(s.shapes[49], "Target")
    set_text(s.shapes[51], "Performance is interpreted in context, never in isolation.")
    set_text(s.shapes[53], "Best practice: read Performance together with Decisions + Runs + Ledger.")

    # Slide 16 - Audit trail
    s = prs.slides[15]
    clear_slide_text(s)
    set_text(s.shapes[1], "Audit Trail (Runs)")
    set_text(s.shapes[5], "Run status, duration, as-of, and portfolio scope")
    set_text(s.shapes[8], "Step-level timeline with pass/fail/skip and timing")
    set_text(s.shapes[10], "Error panel with SQLSTATE and query IDs")
    set_text(s.shapes[13], "Daily + intraday pipeline mode toggle")
    set_text(s.shapes[15], "Triage flow: summary -> step details -> root cause")
    set_text(s.shapes[18], "Cross-check with Live Activity and Decisions for incident review")
    set_text(s.shapes[20], "Demo-ready incident flow")
    set_text(s.shapes[23], "1) Open latest failed run")
    set_text(s.shapes[24], "2) Identify failing procedure and query ID")
    set_text(s.shapes[25], "3) Validate impact on decisions/activity")
    set_text(s.shapes[25], "4) Confirm recovery in next successful run")

    # Slide 17 - Demo script
    s = prs.slides[16]
    clear_slide_text(s)
    set_text(s.shapes[1], "Live Demo Script (Recommended Order)")
    set_text(s.shapes[5], "Step 1")
    set_text(s.shapes[6], "Home / Cockpit: confirm freshness and overnight changes")
    set_text(s.shapes[9], "Step 2")
    set_text(s.shapes[10], "Training Status: evidence maturity and horizon outcomes")
    set_text(s.shapes[13], "Step 3")
    set_text(s.shapes[14], "Market Timeline: one symbol chain S -> P -> T")
    set_text(s.shapes[17], "Step 4")
    set_text(s.shapes[18], "Parallel Worlds: policy health + regret")
    set_text(s.shapes[21], "Step 5")
    set_text(s.shapes[22], "News Intelligence: context, impact, HOT")
    set_text(s.shapes[25], "Step 6")
    set_text(s.shapes[26], "Live Link + Live Activity: controls and committee checkpoint")
    set_text(s.shapes[29], "Step 7")
    set_text(s.shapes[30], "AI Decisions + Symbol Tracker commentary")
    set_text(s.shapes[33], "Step 8")
    set_text(s.shapes[34], "Performance + Runs + Learning Ledger closure")
    set_text(s.shapes[37], "Presenter Tip")
    set_text(s.shapes[38], "Keep one symbol + one portfolio as the story anchor.")
    set_text(s.shapes[41], "Presenter Tip")
    set_text(s.shapes[42], "When AI appears, state guardrails and deterministic gates.")
    set_text(s.shapes[45], "Outcome")
    set_text(s.shapes[46], "Audience leaves with both product and operator understanding.")

    # Slide 18 - Quick reference
    s = prs.slides[17]
    clear_slide_text(s)
    set_text(s.shapes[1], "Quick Reference for Training")
    set_text(s.shapes[3], "Updated terms and checkpoints used during this demo.")
    set_text(s.shapes[6], "Term")
    set_text(s.shapes[9], "Meaning")
    set_text(s.shapes[13], "Committee verdict")
    set_text(s.shapes[14], "Final decision posture for proposal")
    set_text(s.shapes[18], "Reason codes")
    set_text(s.shapes[19], "Structured explanation of blocks/approvals")
    set_text(s.shapes[23], "Thesis")
    set_text(s.shapes[24], "Symbol validity state: intact/weaken/invalid")
    set_text(s.shapes[28], "Open R")
    set_text(s.shapes[29], "Current reward-to-risk multiple")
    set_text(s.shapes[33], "Policy Health")
    set_text(s.shapes[34], "Parallel Worlds quality signal for rules")
    set_text(s.shapes[38], "HOT")
    set_text(s.shapes[39], "Unseen decision-relevant news marker")
    set_text(s.shapes[43], "Runs (Audit)")
    set_text(s.shapes[44], "Operational truth source for pipelines")
    set_text(s.shapes[48], "Learning Ledger")
    set_text(s.shapes[49], "Decision-to-outcome causality trace")
    set_text(s.shapes[75], "Operator principle: trust evidence, verify in audit, then act.")

    # Slide 19 - Thank you
    s = prs.slides[18]
    clear_slide_text(s)
    set_text(s.shapes[23], "Thank You")
    set_text(s.shapes[25], "Q&A + Live Walkthrough")
    set_text(
        s.shapes[26],
        "Deep-dive options: committee logic, live workflow controls, or audit/learning trace.",
    )
    set_text(s.shapes[29], "Market Intelligence Platform  |  Internal Training  |  March 2026")


def build_presentation() -> Path:
    root = Path(__file__).resolve().parents[2]
    template_path = root / "MIP_Presentation_v2.pptx"
    output_path = Path(__file__).with_name("MIP_Demo_Training_Deck_Reference_Style.pptx")

    prs = Presentation(str(template_path))

    apply_content(prs)
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    out = build_presentation()
    print(f"Created: {out}")
