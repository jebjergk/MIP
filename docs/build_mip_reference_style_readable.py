from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Reference style palette (from existing MIP deck family)
NAVY = RGBColor(20, 30, 80)         # #141E50
NAVY_DARK = RGBColor(13, 21, 64)    # #0D1540
CYAN = RGBColor(79, 195, 247)       # #4FC3F7
BLUE = RGBColor(13, 110, 253)       # #0D6EFD
PURPLE = RGBColor(111, 66, 193)     # #6F42C1
GREEN = RGBColor(25, 135, 84)       # #198754
ORANGE = RGBColor(253, 126, 20)     # #FD7E14
RED = RGBColor(220, 53, 69)         # #DC3545
BG = RGBColor(248, 249, 250)        # #F8F9FA
TEXT = RGBColor(26, 26, 46)         # #1A1A2E
MUTED = RGBColor(108, 117, 125)     # #6C757D


def _header(slide, title: str, subtitle: str, section: str = "MIP"):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.85))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.82), Inches(13.333), Inches(0.07))
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.line.fill.background()

    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.2), Inches(0.2), Inches(1.8), Inches(0.42))
    chip.fill.solid()
    chip.fill.fore_color.rgb = NAVY_DARK
    chip.line.color.rgb = CYAN
    tf = chip.text_frame
    tf.text = section
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    t = slide.shapes.add_textbox(Inches(0.7), Inches(1.05), Inches(10.8), Inches(0.8))
    ttf = t.text_frame
    ttf.text = title
    ttf.paragraphs[0].font.size = Pt(32)
    ttf.paragraphs[0].font.bold = True
    ttf.paragraphs[0].font.color.rgb = NAVY_DARK

    s = slide.shapes.add_textbox(Inches(0.7), Inches(1.75), Inches(12), Inches(0.45))
    stf = s.text_frame
    stf.text = subtitle
    stf.paragraphs[0].font.size = Pt(14)
    stf.paragraphs[0].font.color.rgb = MUTED


def _footer(slide, text: str = "MIP | Evidence-first training walkthrough"):
    f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
    f.fill.solid()
    f.fill.fore_color.rgb = RGBColor(233, 236, 239)
    f.line.fill.background()
    tf = slide.shapes.add_textbox(Inches(0.65), Inches(7.24), Inches(11), Inches(0.2)).text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = MUTED


def _card(slide, x, y, w, h, title, bullets, edge=BLUE):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.fill.solid()
    c.fill.fore_color.rgb = RGBColor(255, 255, 255)
    c.line.color.rgb = RGBColor(233, 236, 239)

    e = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    e.fill.solid()
    e.fill.fore_color.rgb = edge
    e.line.fill.background()

    tf = c.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = NAVY_DARK
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"- {b}"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT


def _flow(slide, labels, y=2.4):
    x = 0.7
    w = 1.9
    for i, lbl in enumerate(labels):
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.9))
        b.fill.solid()
        b.fill.fore_color.rgb = RGBColor(255, 255, 255)
        b.line.color.rgb = RGBColor(206, 212, 218)
        tf = b.text_frame
        tf.text = lbl
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = NAVY_DARK
        if i < len(labels) - 1:
            a = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + w + 0.07), Inches(y + 0.2), Inches(0.2), Inches(0.45))
            a.fill.solid()
            a.fill.fore_color.rgb = CYAN
            a.line.fill.background()
        x += 2.2


def _build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def new():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = BG
        return s

    # 1 Title
    s = new()
    _header(s, "Market Intelligence Platform", "Demo + training: research internals to trading operations", "OVERVIEW")
    _card(s, 0.7, 2.45, 6.2, 3.6, "Why MIP", [
        "Evidence-first decision support.",
        "AI explanations with deterministic controls.",
        "Complete signal-to-outcome audit trail.",
    ], edge=CYAN)
    _card(s, 7.1, 2.45, 5.5, 3.6, "Session Goal", [
        "Teach how internals become trading actions.",
        "Show committee reasoning with examples.",
        "End with performance and audit confidence.",
    ], edge=PURPLE)
    _footer(s)

    # 2 Intro
    s = new()
    _header(s, "MIP Introduction", "From research learning to controlled execution", "INTRO")
    _flow(s, ["Market Data", "Signals", "Outcomes", "Trust", "Committee", "Operations"])
    _card(s, 0.7, 3.55, 6.0, 2.8, "Research Functionality", [
        "Pattern detection and horizon evaluation (H1-H20).",
        "Maturity, coverage, and trust scoring.",
        "Counterfactual policy testing (Parallel Worlds).",
    ], edge=BLUE)
    _card(s, 6.9, 3.55, 5.7, 2.8, "Trading Functionality", [
        "Proposals through risk gates and committee checks.",
        "Live-linked activity and symbol monitoring.",
        "Performance + audit + learning feedback loop.",
    ], edge=GREEN)
    _footer(s)

    # 3 Training
    s = new()
    _header(s, "Training Status", "Evidence quality before action", "RESEARCH")
    _card(s, 0.7, 2.3, 6.0, 2.0, "How to read", [
        "Start with sample size + coverage.",
        "Then assess horizon averages and consistency.",
        "Maturity is confidence in evidence, not certainty.",
    ], edge=BLUE)
    _card(s, 6.9, 2.3, 5.7, 2.0, "Example", [
        "AUD/USD Pattern 2",
        "Coverage 86%, Maturity 78",
        "Avg H5 +0.8% (historical edge)",
    ], edge=GREEN)
    _card(s, 0.7, 4.55, 11.9, 1.8, "AI naturally fused", [
        "Ask MIP translates training metrics into plain-language guidance for operators.",
    ], edge=PURPLE)
    _footer(s)

    # 4 Market timeline
    s = new()
    _header(s, "Market Timeline", "Symbol story from signal to trade", "RESEARCH")
    _flow(s, ["Signals (S)", "Proposals (P)", "Trades (T)"], y=2.45)
    _card(s, 0.7, 3.75, 6.0, 2.6, "Visuals to explain", [
        "Tile colors: executed / proposed / signal-only.",
        "Chart overlays for signal, proposal, execution.",
        "Signal-chain tree by portfolio.",
    ], edge=BLUE)
    _card(s, 6.9, 3.75, 5.7, 2.6, "Training angle", [
        "Use one symbol as anchor across slides.",
        "Narrate why flow advanced or stalled.",
        "Link directly to committee decisions.",
    ], edge=CYAN)
    _footer(s)

    # 5 Parallel worlds
    s = new()
    _header(s, "Parallel Worlds", "Counterfactual lab for policy quality", "RESEARCH")
    _card(s, 0.7, 2.3, 4.0, 2.1, "Policy Health", [
        "Badge: Healthy / Watch / Needs Attention",
        "Stability score and regret driver",
    ], edge=GREEN)
    _card(s, 4.9, 2.3, 3.7, 2.1, "Scenarios", [
        "Signal filter",
        "Position size",
        "Entry timing",
        "Baseline (stay cash)",
    ], edge=BLUE)
    _card(s, 8.8, 2.3, 3.8, 2.1, "Confidence", [
        "Strong / Emerging / Weak / Noise",
        "Only strong patterns move to review",
    ], edge=ORANGE)
    _card(s, 0.7, 4.65, 11.9, 1.7, "AI naturally fused", [
        "AI narrative explains divergence and regret trend; changes remain human-governed.",
    ], edge=PURPLE)
    _footer(s)

    # 6 News intelligence
    s = new()
    _header(s, "News Intelligence", "Evidence-backed context with guardrails", "RESEARCH")
    _card(s, 0.7, 2.3, 3.9, 2.9, "Context KPIs", [
        "Symbols with news",
        "HOT symbols",
        "Snapshot freshness",
    ], edge=BLUE)
    _card(s, 4.8, 2.3, 3.9, 2.9, "Decision Impact", [
        "Proposals with news_context",
        "Score adjustments",
        "Blocked new entries",
    ], edge=PURPLE)
    _card(s, 8.9, 2.3, 3.7, 2.9, "Guardrails", [
        "Invalid URLs excluded",
        "Influence bounded",
        "Freshness always explicit",
    ], edge=GREEN)
    _card(s, 0.7, 5.45, 11.9, 0.95, "AI naturally fused", [
        "Reader summary is grounded in stored features, not free-form speculation.",
    ], edge=CYAN)
    _footer(s)

    # 7 Live portfolio link
    s = new()
    _header(s, "Live Portfolio Link", "Control plane from source to broker truth", "TRADING")
    _flow(s, ["Source Portfolio", "Live Portfolio", "IBKR", "Activation Guard", "Readiness"], y=2.45)
    _card(s, 0.7, 3.8, 6.0, 2.55, "What this page controls", [
        "Adapter mode, exposure limits, freshness controls.",
        "Drawdown/bust safety brakes.",
        "Saved config = governance record only.",
    ], edge=BLUE)
    _card(s, 6.9, 3.8, 5.7, 2.55, "What it does not do", [
        "Does not place orders by itself.",
        "Execution still requires approvals and revalidation.",
        "Research source selected during activity import.",
    ], edge=RED)
    _footer(s)

    # 8 Live activity
    s = new()
    _header(s, "Live Portfolio Activity", "Operational lifecycle with committee checkpoint", "TRADING")
    _flow(s, ["Imported", "Validated", "Committee", "Approved", "Executed"], y=2.45)
    _card(s, 0.7, 3.8, 6.0, 2.55, "Use during demo", [
        "Verify freshest transitions and statuses.",
        "Investigate delays with reason fields.",
        "Trace one action from import to execution.",
    ], edge=BLUE)
    _card(s, 6.9, 3.8, 5.7, 2.55, "Committee emphasis", [
        "Decision checkpoint is visible and auditable.",
        "Cross-check with AI Agent Decisions and Runs.",
        "Great training view for process discipline.",
    ], edge=PURPLE)
    _footer(s)

    # 9 AI decisions
    s = new()
    _header(s, "AI Agent Decisions", "Committee courtroom: verdict, reasons, revalidation", "DECISIONS")
    _card(s, 0.7, 2.3, 5.9, 2.0, "Accepted example", [
        "Verdict: APPROVE",
        "Reason: trust + risk + freshness passed",
        "Status path: PROPOSED -> APPROVED -> EXECUTED",
    ], edge=GREEN)
    _card(s, 6.8, 2.3, 5.8, 2.0, "Rejected example", [
        "Verdict: REJECT",
        "Reason: capacity/risk or stale data",
        "Use reason tags for policy tuning",
    ], edge=RED)
    _card(s, 0.7, 4.55, 11.9, 1.8, "AI naturally fused", [
        "Committee summaries accelerate understanding while structured reason codes preserve deterministic accountability.",
    ], edge=PURPLE)
    _footer(s)

    # 10 Symbol tracker
    s = new()
    _header(s, "Live Symbol Tracker", "Live symbol monitoring + committee commentary", "DECISIONS")
    _card(s, 0.7, 2.3, 6.0, 2.0, "Symbol metrics", [
        "Thesis state: intact / weakening / invalidated.",
        "Open R, expected move reached, distance to TP/SL.",
        "Protection and status badges by position.",
    ], edge=BLUE)
    _card(s, 6.9, 2.3, 5.7, 2.0, "Committee commentary", [
        "Stance, confidence, reason tags, actions to consider.",
        "Contextual guidance, non-binding by design.",
        "Escalates symbols needing immediate review.",
    ], edge=PURPLE)
    _card(s, 0.7, 4.55, 11.9, 1.8, "Training angle", [
        "Teach how to combine symbol commentary with AI decisions, news context, and run health.",
    ], edge=CYAN)
    _footer(s)

    # 11 Learning ledger
    s = new()
    _header(s, "Learning Ledger", "Why decisions changed and what happened next", "REVIEW")
    _flow(s, ["Proposal", "Decision", "Execution", "Outcome", "Lesson"], y=2.45)
    _card(s, 0.7, 3.8, 6.0, 2.55, "What it proves", [
        "Evidence chain behind behavior changes.",
        "Expected vs realized outcomes.",
        "Impact of repeated decision patterns.",
    ], edge=BLUE)
    _card(s, 6.9, 3.8, 5.7, 2.55, "How to use it", [
        "Post-trade review and weekly retrospectives.",
        "Stakeholder explainability and audit support.",
        "Policy tuning based on actual outcomes.",
    ], edge=GREEN)
    _footer(s)

    # 12 Performance
    s = new()
    _header(s, "Performance Dashboard", "Portfolio-level truth on return, drawdown, consistency", "REVIEW")
    _card(s, 0.7, 2.3, 3.8, 2.4, "Portfolio A", ["Return +12.4%", "Max DD -4.1%", "Stability High"], edge=GREEN)
    _card(s, 4.75, 2.3, 3.8, 2.4, "Portfolio B", ["Return +8.7%", "Max DD -2.9%", "Stability Very High"], edge=BLUE)
    _card(s, 8.8, 2.3, 3.8, 2.4, "Portfolio C", ["Return +15.1%", "Max DD -7.3%", "Stability Medium"], edge=ORANGE)
    _card(s, 0.7, 4.95, 11.9, 1.4, "Interpretation workflow", [
        "Read period trend first, then connect moves back to training, committee decisions, and risk gates.",
    ], edge=PURPLE)
    _footer(s)

    # 13 Audit trail
    s = new()
    _header(s, "Audit Trail (Runs)", "Operational truth source for daily and intraday pipelines", "OPERATIONS")
    _card(s, 0.7, 2.3, 5.9, 2.9, "Runs panel", [
        "Status, timing, as-of, and portfolio scope.",
        "Step timeline with pass/fail/skip details.",
        "Error diagnostics: SQLSTATE + query ID.",
    ], edge=BLUE)
    _card(s, 6.8, 2.3, 5.8, 2.9, "Incident playbook", [
        "1) Open failed run",
        "2) Locate failing step and query id",
        "3) Validate downstream impact",
        "4) Confirm recovery on next success",
    ], edge=RED)
    _card(s, 0.7, 5.45, 11.9, 0.95, "AI naturally fused", [
        "AI run summary speeds triage, but step-level diagnostics remain source of truth.",
    ], edge=CYAN)
    _footer(s)

    # 14 AI fusion
    s = new()
    _header(s, "How AI Is Naturally Fused", "AI improves speed and clarity while controls preserve safety", "GOVERNANCE")
    _card(s, 0.7, 2.3, 5.9, 4.1, "Where AI helps", [
        "Committee summaries and narrative context.",
        "Parallel Worlds explanation and regret interpretation.",
        "News reader summary and HOT prioritization.",
        "Ask MIP route-aware operator coaching.",
    ], edge=PURPLE)
    _card(s, 6.8, 2.3, 5.8, 4.1, "Where humans stay in control", [
        "Deterministic risk gates and thresholds.",
        "Execution approvals and revalidation.",
        "AI outputs are non-binding.",
        "Runs + Ledger ensure accountability.",
    ], edge=GREEN)
    _footer(s)

    # 15 Demo script
    s = new()
    _header(s, "Recommended Demo Script", "Use this sequence tomorrow for product + training impact", "SCRIPT")
    _card(s, 0.7, 2.3, 11.9, 3.9, "8-step run order", [
        "1) Home/Cockpit: freshness and overnight changes.",
        "2) Training Status: evidence maturity and horizons.",
        "3) Market Timeline: one symbol chain S -> P -> T.",
        "4) Parallel Worlds: policy health and confidence.",
        "5) News Intelligence: context + impact + HOT.",
        "6) Live Link + Live Activity: controls and lifecycle.",
        "7) AI Decisions + Symbol Tracker commentary.",
        "8) Performance + Runs + Learning Ledger close.",
    ], edge=CYAN)
    _footer(s)

    # 16 Close
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    title = s.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.8), Inches(1.0)).text_frame
    title.text = "Q&A"
    title.paragraphs[0].font.size = Pt(56)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    sub = s.shapes.add_textbox(Inches(0.95), Inches(2.8), Inches(11.0), Inches(2.6)).text_frame
    sub.text = "Deep dive options:"
    sub.paragraphs[0].font.size = Pt(20)
    sub.paragraphs[0].font.color.rgb = RGBColor(221, 238, 255)
    for t in [
        "Committee logic and rejection patterns",
        "Live workflow controls and readiness",
        "Audit + learning trace from decision to outcome",
    ]:
        p = sub.add_paragraph()
        p.text = f"- {t}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)

    out = Path(__file__).with_name("MIP_Demo_Training_Deck_Reference_Style_Readable.pptx")
    prs.save(out)
    print(f"Created: {out}")


if __name__ == "__main__":
    _build()
